#!/usr/bin/env python3
"""AI Manga Translator 初步验收汇报 PPT 构建脚本。

风格：瑞士国际主义（Swiss International）
- 16:9，深墨黑/纸白底 + 单一克莱因蓝（IKB #002FA7）强调色
- 无衬线字体、左对齐网格、深色封面/章节页与浅色内容页交替
- 全部元素为原生可编辑形状/文本框/表格；图片位为带边框占位框

运行：.venv\\Scripts\\python.exe docs\\正式交付文档\\验收阶段\\ppt\\build_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "AI_Manga_Translator_验收汇报.pptx"

PAPER = RGBColor(0xFA, 0xFA, 0xF8)
INK = RGBColor(0x0A, 0x0A, 0x0A)
ACCENT = RGBColor(0x00, 0x2F, 0xA7)
GREY_1 = RGBColor(0xF0, 0xF0, 0xEE)
GREY_2 = RGBColor(0xD4, 0xD4, 0xD2)
GREY_3 = RGBColor(0x73, 0x73, 0x73)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_DIM = RGBColor(0xC9, 0xC9, 0xC7)

F_TITLE = "Microsoft YaHei Light"
F_BODY = "Microsoft YaHei"
F_MONO = "Consolas"

PAGE_W = 13.333
PAGE_H = 7.5
MARGIN = 0.6
CONTENT_W = PAGE_W - MARGIN * 2
FOOT_Y = 7.02
TOTAL_PAGES = 15


def _set_run(run, text, size, color, bold=False, font=F_BODY, italic=False, spacing=None):
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def tb(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dict(text, size, color, bold, font, align, italic, space_before, space_after, line_spacing, spacing)"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if spec.get("space_before") is not None:
            p.space_before = Pt(spec["space_before"])
        if spec.get("space_after") is not None:
            p.space_after = Pt(spec["space_after"])
        if spec.get("line_spacing") is not None:
            p.line_spacing = spec["line_spacing"]
        runs = spec.get("runs") or [spec]
        for rs in runs:
            _set_run(
                p.add_run(),
                rs.get("text", ""),
                rs.get("size", spec.get("size", 11)),
                rs.get("color", spec.get("color", INK)),
                rs.get("bold", spec.get("bold", False)),
                rs.get("font", spec.get("font", F_BODY)),
                rs.get("italic", spec.get("italic", False)),
                rs.get("spacing", spec.get("spacing")),
            )
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, dash=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
        if dash is not None:
            ln = sp.line._get_or_add_ln()
            prst = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
            ln.append(prst)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return sp


def shape_text(sp, paras, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.CENTER)
        if spec.get("space_after") is not None:
            p.space_after = Pt(spec["space_after"])
        if spec.get("line_spacing") is not None:
            p.line_spacing = spec["line_spacing"]
        for rs in (spec.get("runs") or [spec]):
            _set_run(
                p.add_run(),
                rs.get("text", ""),
                rs.get("size", spec.get("size", 10)),
                rs.get("color", spec.get("color", INK)),
                rs.get("bold", spec.get("bold", False)),
                rs.get("font", spec.get("font", F_BODY)),
                rs.get("italic", spec.get("italic", False)),
                rs.get("spacing", spec.get("spacing")),
            )


def hline(slide, x, y, w, color=GREY_2, weight=0.75):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def chrome(slide, page_no, section, dark=False):
    c_meta = PAPER_DIM if dark else GREY_3
    tb(slide, MARGIN, 0.32, 8.5, 0.3, [
        {"text": f"AI Manga Translator · {section}", "size": 9, "color": c_meta, "font": F_MONO, "spacing": 1.6},
    ])
    tb(slide, PAGE_W - MARGIN - 2.5, 0.32, 2.5, 0.3, [
        {"text": f"{page_no:02d} / {TOTAL_PAGES}", "size": 9, "color": c_meta, "font": F_MONO, "align": PP_ALIGN.RIGHT, "spacing": 1.6},
    ])
    hline(slide, MARGIN, FOOT_Y, CONTENT_W, GREY_2 if not dark else RGBColor(0x3A, 0x3A, 0x3A))
    tb(slide, MARGIN, FOOT_Y + 0.06, 8.5, 0.26, [
        {"text": "漫画多语言智能翻译系统 · 初步验收汇报", "size": 8, "color": c_meta, "font": F_MONO, "spacing": 1.2},
    ])
    tb(slide, PAGE_W - MARGIN - 3.5, FOOT_Y + 0.06, 3.5, 0.26, [
        {"text": f"{section} · {page_no:02d} / {TOTAL_PAGES}", "size": 8, "color": c_meta, "font": F_MONO, "align": PP_ALIGN.RIGHT, "spacing": 1.2},
    ])


def header(slide, kicker, title, conclusion, dark=False, title_size=27):
    kc = RGBColor(0x8F, 0xA8, 0xE8) if dark else ACCENT
    tc = WHITE if dark else INK
    cc = PAPER_DIM if dark else GREY_3
    tb(slide, MARGIN, 0.78, CONTENT_W, 0.28, [
        {"text": kicker, "size": 10, "color": kc, "font": F_MONO, "bold": True, "spacing": 2.2},
    ])
    tb(slide, MARGIN, 1.06, CONTENT_W, 0.62, [
        {"text": title, "size": title_size, "color": tc, "font": F_TITLE},
    ])
    tb(slide, MARGIN, 1.66, CONTENT_W, 0.32, [
        {"text": conclusion, "size": 12, "color": cc},
    ])


def placeholder(slide, x, y, w, h, name, desc, dark=False):
    fill = DARK_CARD if dark else WHITE
    line = RGBColor(0x5A, 0x6F, 0xB0) if dark else ACCENT
    label_c = RGBColor(0x9D, 0xB1, 0xE8) if dark else ACCENT
    desc_c = PAPER_DIM if dark else GREY_3
    box = rect(slide, x, y, w, h, fill=fill, line=line, line_w=1.1, dash=True)
    shape_text(box, [
        {"text": name, "size": 10.5, "color": label_c, "bold": True, "font": F_MONO, "space_after": 6},
        {"text": f"图片说明：{desc}", "size": 8.5, "color": desc_c, "line_spacing": 1.25},
    ])
    return box


def card(slide, x, y, w, h, title, body_lines, kind="fill", num=None, title_size=12.5, body_size=10):
    bg = {"fill": GREY_1, "ink": INK, "accent": ACCENT, "outline": None}[kind]
    fg = WHITE if kind in ("ink", "accent") else INK
    sub = PAPER_DIM if kind in ("ink", "accent") else GREY_3
    sp = rect(slide, x, y, w, h, fill=bg, line=(GREY_2 if kind == "outline" else None))
    paras = []
    if num is not None:
        paras.append({"text": num, "size": 9, "color": (ACCENT if kind == "fill" else sub), "font": F_MONO, "bold": True, "align": PP_ALIGN.LEFT, "space_after": 3, "spacing": 1.5})
    paras.append({"text": title, "size": title_size, "color": fg, "bold": True, "align": PP_ALIGN.LEFT, "space_after": 4})
    for ln in body_lines:
        paras.append({"text": ln, "size": body_size, "color": (fg if kind in ("ink", "accent") else RGBColor(0x3D, 0x3D, 0x3D)), "align": PP_ALIGN.LEFT, "line_spacing": 1.22, "space_after": 2})
    shape_text(sp, paras, anchor=MSO_ANCHOR.TOP)
    sp.text_frame.margin_left = Inches(0.14)
    sp.text_frame.margin_right = Inches(0.12)
    sp.text_frame.margin_top = Inches(0.1)
    return sp


def chevron_row(slide, x, y, w, h, items, fill=GREY_1, fg=INK, size=10.5, gap=0.08, accent_last=False, number=True):
    n = len(items)
    cw = (w - gap * (n - 1)) / n
    for i, it in enumerate(items):
        last = accent_last and i == n - 1
        sp = rect(
            slide, x + i * (cw + gap), y, cw, h,
            fill=(ACCENT if last else fill),
            shape=(MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON),
        )
        label = f"{i + 1:02d}  {it}" if number else it
        shape_text(sp, [{"text": label, "size": size, "color": (WHITE if last else fg), "bold": True, "line_spacing": 1.1}])


def arrow_between(slide, x, y, size=0.22, color=GREY_3):
    sp = rect(slide, x, y, size, size * 0.62, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)
    return sp


def make_table(slide, x, y, w, h, headers, rows, col_widths=None, header_fill=INK, font_size=10):
    g = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = g.table
    if col_widths:
        for i, cwr in enumerate(col_widths):
            table.columns[i].width = Inches(cwr)
    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.08)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        p = cell.text_frame.paragraphs[0]
        _set_run(p.add_run(), htxt, font_size, WHITE, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else GREY_1
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            _set_run(p.add_run(), val, font_size, INK)
    return g


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]


def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = INK if dark else PAPER
    return s


# ============================================================
# 第 01 页 · 封面：项目定位（统一浅色风格）
# ============================================================
s = new_slide()
rect(s, 0, 0, 0.16, PAGE_H, fill=ACCENT)
tb(s, 0.9, 0.9, 6.6, 0.3, [
    {"text": "初步验收汇报 · PRELIMINARY ACCEPTANCE REVIEW", "size": 10.5, "color": ACCENT, "font": F_MONO, "bold": True, "spacing": 2.4},
])
tb(s, 0.9, 1.5, 6.9, 2.6, [
    {"text": "AI Manga", "size": 48, "color": INK, "font": F_TITLE, "line_spacing": 1.02},
    {"text": "Translator", "size": 48, "color": INK, "font": F_TITLE, "line_spacing": 1.02},
])
tb(s, 0.9, 3.55, 6.8, 0.5, [
    {"text": "漫画多语言智能翻译系统", "size": 20, "color": INK, "font": F_BODY},
])
tb(s, 0.9, 4.15, 6.6, 0.75, [
    {"text": "面向漫画图片的 OCR、翻译、擦除修复与中文排版一体化工具", "size": 12.5, "color": GREY_3, "line_spacing": 1.35},
])
hline(s, 0.9, 5.35, 6.3, GREY_2)
tb(s, 0.9, 5.5, 6.8, 1.3, [
    {"runs": [
        {"text": "汇报人 ", "size": 10.5, "color": GREY_3, "font": F_MONO},
        {"text": "【姓名】", "size": 10.5, "color": INK},
        {"text": "    学号 ", "size": 10.5, "color": GREY_3, "font": F_MONO},
        {"text": "【学号】", "size": 10.5, "color": INK},
    ], "space_after": 6},
    {"runs": [
        {"text": "班级   ", "size": 10.5, "color": GREY_3, "font": F_MONO},
        {"text": "【学习方向或班级】", "size": 10.5, "color": INK},
        {"text": "    指导教师 ", "size": 10.5, "color": GREY_3, "font": F_MONO},
        {"text": "【教师姓名】", "size": 10.5, "color": INK},
    ], "space_after": 6},
    {"runs": [
        {"text": "日期   ", "size": 10.5, "color": GREY_3, "font": F_MONO},
        {"text": "【汇报日期】", "size": 10.5, "color": INK},
        {"text": "    项目地址 ", "size": 10.5, "color": GREY_3, "font": F_MONO},
        {"text": "【如有则填写】", "size": 10.5, "color": INK},
    ]},
])
placeholder(s, 8.15, 1.5, 4.55, 2.56, "[图片占位-01-项目主界面或原译对照-16:9]",
            "放入实际项目首页或一组真实漫画原图/中文译图，用于建立项目第一印象。")
tb(s, 8.15, 4.35, 4.55, 1.9, [
    {"text": "本地 Web 应用", "size": 11, "color": INK, "bold": True, "space_after": 4},
    {"text": "React/Vite 前端 + FastAPI 后端", "size": 10, "color": GREY_3, "space_after": 2, "line_spacing": 1.3},
    {"text": "SQLite 与本地文件存储", "size": 10, "color": GREY_3, "space_after": 2, "line_spacing": 1.3},
    {"text": "图片输入 · 日语/英语 → 中文", "size": 10, "color": GREY_3, "line_spacing": 1.3},
])
hline(s, 0.9, FOOT_Y, PAGE_W - 0.9 - MARGIN, GREY_2)
tb(s, 0.9, FOOT_Y + 0.06, 8.5, 0.26, [
    {"text": "AI MANGA TRANSLATOR · COVER", "size": 8, "color": GREY_3, "font": F_MONO, "spacing": 1.2},
])
tb(s, PAGE_W - MARGIN - 3.5, FOOT_Y + 0.06, 3.5, 0.26, [
    {"text": "01 / 15", "size": 8, "color": GREY_3, "font": F_MONO, "align": PP_ALIGN.RIGHT, "spacing": 1.2},
])
notes(s, "开场先给定位：这是一个本地运行的漫画图片翻译 Web 应用，输入漫画图片，输出已排版的中文译图。"
         "右侧占位图在汇报前替换为真实首页或原图/译图对照，姓名、学号、班级、指导教师也在此页补全。")

# ============================================================
# 第 02 页 · 为什么要做：问题与价值
# ============================================================
s = new_slide()
chrome(s, 2, "问题与价值")
header(s, "WHY · 02", "为什么要做：问题与价值",
       "漫画翻译的难点不在“识字”，而在擦除、排版与气泡上下文的完整闭环")
card(s, MARGIN, 2.2, 2.6, 1.18, "文字形态复杂", ["竖排、弯曲、密集排列", "文字叠在复杂背景上"], num="痛点 01")
card(s, MARGIN + 2.78, 2.2, 2.6, 1.18, "普通 OCR 只管识字", ["只输出文本，不擦除原文", "不修复背景、不重排译文"], num="痛点 02")
card(s, MARGIN + 5.56, 2.2, 2.6, 1.18, "逐句翻译丢上下文", ["孤立句子丢失气泡语境", "专有名词前后不一致"], num="痛点 03")
tb(s, MARGIN, 4.12, 8.2, 0.26, [{"text": "原始处理链（人工）", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "spacing": 1.5}])
chevron_row(s, MARGIN, 4.42, 8.2, 0.52, ["截图", "识字", "翻译", "修图", "排版"], size=10)
tb(s, MARGIN, 5.12, 8.2, 0.26, [{"text": "本项目处理链（自动）", "size": 10, "color": ACCENT, "font": F_MONO, "bold": True, "spacing": 1.5}])
chevron_row(s, MARGIN, 5.42, 8.2, 0.52, ["上传图片", "自动任务流", "预览编辑", "下载结果"],
            fill=INK, fg=WHITE, size=10)
placeholder(s, 9.35, 2.2, 3.38, 2.25, "[图片占位-02-复杂漫画文字局部-3:2]",
            "放入真实日文或英文漫画原图局部，标出竖排文字、气泡边界和复杂背景。")
tb(s, 9.35, 4.7, 3.4, 2.0, [
    {"text": "标注建议", "size": 10.5, "color": INK, "bold": True, "space_after": 4},
    {"text": "· 竖排文字列", "size": 9.5, "color": GREY_3, "space_after": 2},
    {"text": "· 气泡边界轮廓", "size": 9.5, "color": GREY_3, "space_after": 2},
    {"text": "· 复杂背景上的文字", "size": 9.5, "color": GREY_3, "line_spacing": 1.3},
])
notes(s, "从三个真实痛点讲起：竖排密集文字、普通 OCR 不管擦除排版、逐句翻译丢上下文。"
         "然后指向对照链：过去人工五步，现在上传后自动完成。右侧原图局部用于让难点可视化。")

# ============================================================
# 第 03 页 · 做什么：需求范围与实现边界
# ============================================================
s = new_slide()
chrome(s, 3, "范围与边界")
header(s, "SCOPE · 03", "做什么：需求范围与实现边界",
       "本次验收只针对“当前已实现”列；依赖环境与后续规划单独标明，不混入")
col_w = (CONTENT_W - 0.6) / 3
headers3 = [("当前已实现", INK), ("依赖环境", ACCENT), ("后续规划", GREY_3)]
body3 = [
    ["图片 / 文件夹输入", "日语、英语 → 中文重点翻译", "任务进度与结果预览", "单张重试", "单张 / 批量 ZIP 下载", "专有名词库管理", "文本框编辑", "局部擦除与撤销"],
    ["真实 OCR 模型：取决于本地安装与模型权重", "LaMa 修复：需额外配置", "外部翻译 API：取决于密钥与网络", "核心 Demo 模式在较少依赖下可运行"],
    ["历史翻译记录", "章节级管理", "EPUB / CBZ / PDF", "团队协作", "云端 GPU"],
]
for i, (label, c) in enumerate(headers3):
    x = MARGIN + i * (col_w + 0.3)
    rect(s, x, 2.25, col_w, 0.5, fill=c)
    tb(s, x + 0.12, 2.33, col_w - 0.2, 0.34, [{"text": label, "size": 12, "color": WHITE, "bold": True}])
    box = rect(s, x, 2.75, col_w, 3.2, fill=(GREY_1 if i != 2 else WHITE), line=(GREY_2 if i == 2 else None))
    paras = []
    for ln in body3[i]:
        paras.append({"text": "· " + ln, "size": 10, "color": RGBColor(0x3D, 0x3D, 0x3D), "align": PP_ALIGN.LEFT, "line_spacing": 1.25, "space_after": 5})
    shape_text(box, paras, anchor=MSO_ANCHOR.TOP)
    box.text_frame.margin_left = Inches(0.14)
    box.text_frame.margin_top = Inches(0.12)
tb(s, MARGIN, 6.6, CONTENT_W, 0.3, [
    {"text": "边界声明：规划列内容不计入本次验收；依赖列能力是否可用以现场环境为准。", "size": 9.5, "color": GREY_3, "font": F_MONO},
])
notes(s, "这一页主动划边界：左列才是本次验收范围；中列说明真实 OCR 模型、LaMa、翻译 API 取决于本机安装、权重、密钥和网络；"
         "右列是规划，不要当成已实现。被追问进度时回到这三栏回答。")

# ============================================================
# 第 04 页 · 用户怎么用：一次完整任务
# ============================================================
s = new_slide()
chrome(s, 4, "用户任务流")
header(s, "USER FLOW · 04", "用户怎么用：一次完整任务",
       "从上传到下载是一条任务流，每一步都有可观察结果")
steps4 = ["上传/粘贴图片或文件夹", "选择语言与翻译设置", "提交任务", "查看每张图进度", "查看完成预览", "编辑或重新翻译", "下载结果"]
obs4 = ["任务卡出现在侧边栏", "源/目标语言即时生效", "后端校验并创建任务", "进度条逐张推进", "单张完成立即可预览", "失败子任务单独重试", "单张下载或批量 ZIP"]
n4 = len(steps4)
gap4 = 0.07
cw4 = (CONTENT_W - gap4 * (n4 - 1)) / n4
for i, (st, ob) in enumerate(zip(steps4, obs4)):
    x = MARGIN + i * (cw4 + gap4)
    sp = rect(s, x, 2.2, cw4, 0.62, fill=(ACCENT if i in (0, n4 - 1) else GREY_1),
              shape=(MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON))
    shape_text(sp, [{"text": f"{i + 1:02d}", "size": 8.5, "color": (WHITE if i in (0, n4 - 1) else GREY_3), "font": F_MONO, "bold": True, "space_after": 2},
                    {"text": st, "size": 9, "color": (WHITE if i in (0, n4 - 1) else INK), "bold": True, "line_spacing": 1.1}])
    tb(s, x, 2.95, cw4 + 0.12, 1.0, [
        {"text": ob, "size": 8.5, "color": GREY_3, "line_spacing": 1.25},
    ])
hline(s, MARGIN, 4.12, CONTENT_W, GREY_2)
card(s, MARGIN, 4.4, 5.5, 2.35, "任务流形态", [
    "主区域为接近对话式的任务流界面",
    "每张图片是独立子任务，互不影响",
    "失败子任务可单独重新翻译",
    "单张下载与批量 ZIP 下载并存",
], num="可观察结果")
placeholder(s, 6.7, 4.4, 6.03, 2.39, "[图片占位-04-任务流界面录屏截图-16:9]",
            "放入上传、任务进度、结果预览中的 2—3 张真实界面截图。")
notes(s, "按七个真实操作顺序讲一遍任务流，强调每步都有可观察结果：任务卡出现、逐张进度、单张完成即可预览、失败可单独重试。"
         "右下角放真实界面截图或录屏截帧，讲解时指向对应入口。")

# ============================================================
# 第 05 页 · 系统怎么组成：总体架构
# ============================================================
s = new_slide()
chrome(s, 5, "总体架构")
header(s, "ARCHITECTURE · 05", "系统怎么组成：总体架构",
       "五层结构自上而下；批量只是编排，模型推理始终串行")
layers = [
    ("L1", "React/Vite 前端交互层", "任务创建、进度展示、结果预览与编辑交互"),
    ("L2", "FastAPI API 层", "请求校验、任务状态与下载接口"),
    ("L3", "TranslationTaskManager 任务层", "单 worker 调度，流水线锁保护推理"),
    ("L4", "Pipeline 推理层", "检测→OCR→修复→分组→翻译→检查→渲染"),
    ("L5", "SQLite / 本地文件存储层", "任务元数据与结果路径"),
]
ly = 2.2
for i, (nb, name, duty) in enumerate(layers):
    fill = INK if i == 2 else (GREY_1 if i % 2 == 0 else WHITE)
    fg = WHITE if fill == INK else INK
    sub = PAPER_DIM if fill == INK else GREY_3
    sp = rect(s, MARGIN, ly, 8.0, 0.78, fill=fill, line=(GREY_2 if fill == WHITE else None))
    shape_text(sp, [
        {"runs": [
            {"text": nb + "  ", "size": 10, "color": (RGBColor(0x8F, 0xA8, 0xE8) if fill == INK else ACCENT), "font": F_MONO, "bold": True},
            {"text": name, "size": 12, "color": fg, "bold": True},
            {"text": "    " + duty, "size": 9.5, "color": sub},
        ], "align": PP_ALIGN.LEFT},
    ])
    sp.text_frame.margin_left = Inches(0.16)
    if i < len(layers) - 1:
        tb(s, MARGIN + 0.35, ly + 0.74, 0.5, 0.22, [{"text": "▼", "size": 8, "color": GREY_3}])
    ly += 0.93
card(s, 9.0, 2.2, 3.73, 1.98, "串行与批量", [
    "批量任务 = 编排多个子任务",
    "模型推理 = 单 worker 串行",
    "流水线锁做双重保护",
], kind="accent", num="关键约束")
tb(s, 9.0, 4.32, 3.73, 0.75, [
    {"text": "实线 = 数据流（自上而下）", "size": 9, "color": GREY_3, "space_after": 2},
    {"text": "虚线 = 配置 / 引擎回退关系", "size": 9, "color": GREY_3},
])
placeholder(s, 9.0, 5.15, 3.73, 1.68, "[图片占位-05-架构图或代码结构-16:9]",
            "优先放实际架构图或后端/前端目录截图，证明模块不是概念拼接。")
notes(s, "自上而下讲五层职责，重点停在第三层：任务管理器是单 worker，加流水线锁。"
         "批量只是编排多个子任务，模型推理仍然串行——这是稳定性优先的取舍，不要讲成并发加速。")

# ============================================================
# 第 06 页 · 核心流程（统一浅色风格）
# ============================================================
s = new_slide()
chrome(s, 6, "核心流水线")
header(s, "PIPELINE · 06", "核心流程：一张图如何变成译图",
       "先修复背景，再按气泡分组整块翻译，最后受控渲染")
flow6 = [
    ("detect", "检测文字区域", "入：原图 / 出：文字框"),
    ("OCR", "识别文本", "入：文字区域 / 出：文本+置信度+方向"),
    ("inpaint", "擦除修复", "入：掩膜 / 出：干净背景"),
    ("bubble grouping", "气泡分组", "入：清理图+原图 / 出：气泡组"),
    ("translate", "整块翻译", "入：气泡文本 / 出：译文"),
    ("quality", "质量检查", "入：译文 / 出：告警标记"),
    ("render", "排版渲染", "入：译文+掩膜 / 出：译图"),
]
n6 = len(flow6)
gap6 = 0.1
cw6 = (CONTENT_W - gap6 * (n6 - 1)) / n6
for i, (en, zh, io) in enumerate(flow6):
    x = MARGIN + i * (cw6 + gap6)
    sp = rect(s, x, 2.25, cw6, 1.5, fill=GREY_1, line=GREY_2, line_w=0.75)
    shape_text(sp, [
        {"text": f"{i + 1:02d}", "size": 9, "color": ACCENT, "font": F_MONO, "bold": True, "space_after": 3},
        {"text": zh, "size": 10.5, "color": INK, "bold": True, "space_after": 2, "line_spacing": 1.1},
        {"text": en, "size": 8, "color": GREY_3, "font": F_MONO, "space_after": 4},
        {"text": io, "size": 7.5, "color": GREY_3, "line_spacing": 1.2},
    ], anchor=MSO_ANCHOR.TOP)
    sp.text_frame.margin_top = Inches(0.09)
notes6 = [
    ("检测", "原图保留边界证据，清理图只用于判断气泡内部连通性，两者分工不混。"),
    ("翻译", "翻译单位是气泡整块，不是孤立字符；提示词要求保留原文换行与句子顺序。"),
    ("渲染", "译文按气泡掩膜裁剪；几何不可靠时按安全扩展框、紧致框分级回退。"),
]
for i, (t, d) in enumerate(notes6):
    x = MARGIN + i * 2.68
    card(s, x, 4.05, 2.52, 1.55, t, [d], kind="fill", num=f"判断 0{i + 1}", title_size=11.5, body_size=8.5)
placeholder(s, 8.62, 4.05, 4.11, 2.31, "[图片占位-06-流水线实测流程图或调试截图-16:9]",
            "可放代码流程截图、任务日志截图或自行绘制的真实流程示意。")
tb(s, MARGIN, 5.85, 8.0, 0.6, [
    {"text": "为什么先修复再分组翻译：分组在干净背景上做连通性判断，避免原文笔画把相邻气泡错误连通；翻译按气泡整块进行，保住上下文。",
     "size": 10, "color": GREY_3, "line_spacing": 1.35},
])
notes(s, "这一页是技术核心，按 detect 到 render 七步讲输入输出。重点解释两个设计：原图保留边界证据、清理图判断连通；"
         "翻译单位是气泡整块。不要逐字念，指着流程图讲输入、处理、输出。")

# ============================================================
# 第 07 页 · 识别难点：OCR、方向与引擎回退
# ============================================================
s = new_slide()
chrome(s, 7, "OCR 与引擎回退")
header(s, "OCR · 07", "识别难点：OCR、方向与引擎回退",
       "检测先定位、识别分主备；manga-ocr 只救空行，不覆盖已读对的结果")
tb(s, MARGIN, 2.2, 4.4, 0.28, [{"text": "识别结果长什么样", "size": 11, "color": INK, "bold": True}])
ocr_flow = [("文字框 / 多边形", "检测定位文字区域"), ("识别文本 + 置信度", "逐区域识别"), ("方向：横排 / 竖排", "供擦除与排版使用")]
oy = 2.56
for i, (t, d) in enumerate(ocr_flow):
    sp = rect(s, MARGIN, oy, 4.4, 0.82, fill=GREY_1)
    shape_text(sp, [
        {"text": t, "size": 11.5, "color": INK, "bold": True, "space_after": 2},
        {"text": d, "size": 9, "color": GREY_3},
    ])
    if i < 2:
        tb(s, MARGIN + 2.0, oy + 0.8, 0.5, 0.2, [{"text": "▼", "size": 8, "color": GREY_3}])
    oy += 1.02
tb(s, MARGIN, oy + 0.06, 4.5, 0.9, [
    {"text": "输出字段：文本、置信度、文字框/多边形、方向；竖排区域旋转后再识别，坐标映射回原图。",
     "size": 9.5, "color": GREY_3, "line_spacing": 1.3},
])
tb(s, 5.5, 2.2, 4.0, 0.28, [{"text": "引擎路由（可插拔，非并行）", "size": 11, "color": INK, "bold": True}])
make_table(s, 5.5, 2.56, 7.2, 2.6,
           ["引擎", "角色", "说明"],
           [["MIT / CTD 检测", "默认路线", "定位文字行并生成掩膜"],
            ["MIT48 识别", "主识别", "输出文本、置信度与方向"],
            ["manga-ocr", "救空行", "仅接管 MIT48 完全未读出的空行"],
            ["PaddleOCR / CV", "回退路径", "可配置启用，不与默认路线并行"]],
           col_widths=[2.1, 1.4, 3.7], font_size=9.5)
placeholder(s, 5.5, 5.4, 4.0, 1.5, "[图片占位-07-OCR原图与检测框/识别结果-3:2]",
            "放入真实检测框截图或 OCR 调试结果，突出竖排和风格化字体。")
tb(s, 9.75, 5.4, 3.0, 1.5, [
    {"text": "横排 / 竖排差异", "size": 10, "color": INK, "bold": True, "space_after": 3},
    {"text": "竖排经旋转后识别；方向信息参与后续分列排版。", "size": 9, "color": GREY_3, "line_spacing": 1.3},
])
notes(s, "先讲识别结果包含什么：文本、置信度、多边形、方向，后续擦除与排版都依赖它们。"
         "再讲引擎路由：MIT48 主识别，manga-ocr 只救完全读不出的空行，PaddleOCR/CV 只是回退路径，不是并行。")

# ============================================================
# 第 08 页 · 图像重建：擦除、修复与气泡分组
# ============================================================
s = new_slide()
chrome(s, 8, "擦除修复与分组")
header(s, "INPAINT & GROUPING · 08", "图像重建：擦除、修复与气泡分组",
       "掩膜合并、分级修复与原图边界否决，共同降低误擦与跨气泡合并")
placeholder(s, MARGIN, 2.2, 5.2, 3.9, "[图片占位-08-原图掩膜修复分组四格对照-4:3]",
            "放入同一张真实图片的四阶段结果：原图 / 掩膜 / 修复图 / 分组结果，证明修复不是简单涂白。")
tb(s, MARGIN, 6.25, 5.2, 0.5, [
    {"text": "四格顺序：原图 → 掩膜 → 修复图 → 分组结果", "size": 9, "color": GREY_3, "font": F_MONO},
])
card(s, 6.2, 2.2, 6.5, 1.05, "掩膜由三部分合并", [
    "文字多边形整块填充 + Otsu 笔画候选 + 竖排注音扩展，覆盖抗锯齿边缘",
], num="机制 01", title_size=11.5, body_size=9.5)
card(s, 6.2, 3.37, 6.5, 1.05, "默认 OpenCV 局部修复，可选 LaMa", [
    "LaMa 需额外配置，是否启用以本机环境为准；刊头/拟声词等非气泡文字保留，不进入擦除、翻译和渲染",
], num="机制 02", title_size=11.5, body_size=9.5)
card(s, 6.2, 4.54, 6.5, 1.05, "清理图分组 + 原图边界否决", [
    "在清理图上做气泡连通分组；两组文字之间若存在原图长轮廓则禁止合并，避免整页吞并",
], num="机制 03", title_size=11.5, body_size=9.5)
tb(s, 6.2, 5.75, 6.5, 0.26, [{"text": "气泡几何回退链", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "spacing": 1.5}])
chevron_row(s, 6.2, 6.03, 6.5, 0.56, ["可靠气泡掩膜", "安全扩展框", "紧致文本框", "跳过该气泡"],
            size=9, number=False, gap=0.06)
notes(s, "用四格对照讲修复不是涂白：掩膜来自文字多边形、Otsu 笔画与注音扩展。强调清理图分组加原图边界否决如何防止跨气泡合并，"
         "最后指回退链：可靠掩膜、安全扩展框、紧致框、跳过。")

# ============================================================
# 第 09 页 · 翻译怎么更稳：上下文、回退与质量检查
# ============================================================
s = new_slide()
chrome(s, 9, "翻译与质量")
header(s, "TRANSLATION · 09", "翻译怎么更稳：上下文、回退与质量检查",
       "以气泡整块翻译保留上下文，用回退链和质量检查兜底")
flow9 = ["气泡内文本按行拼接", "一次整块翻译", "术语库整词替换", "质量检查"]
fy = 2.25
for i, t in enumerate(flow9):
    sp = rect(s, MARGIN, fy, 4.0, 0.66, fill=(GREY_1 if i < 3 else INK))
    shape_text(sp, [{"text": f"{i + 1:02d}  {t}", "size": 10.5, "color": (INK if i < 3 else WHITE), "bold": True, "align": PP_ALIGN.LEFT}])
    sp.text_frame.margin_left = Inches(0.14)
    if i < 3:
        tb(s, MARGIN + 1.8, fy + 0.63, 0.5, 0.2, [{"text": "▼", "size": 8, "color": GREY_3}])
    fy += 0.87
tb(s, MARGIN, fy + 0.08, 4.1, 1.4, [
    {"text": "翻译后端按配置运行，失败进入回退链", "size": 9.5, "color": INK, "bold": True, "space_after": 3},
    {"text": "DeepSeek / OpenAI / DeepL / Google / MyMemory 等适配能力，是否可用取决于本地配置与网络。",
     "size": 9, "color": GREY_3, "line_spacing": 1.3},
])
tb(s, 5.1, 2.2, 4.6, 0.28, [{"text": "质量检查告警项", "size": 11, "color": INK, "bold": True}])
make_table(s, 5.1, 2.56, 7.6, 2.85,
           ["检查项", "判定依据"],
           [["免费/原文回退提示", "实际后端写入记录，回退可见"],
            ["译文明显过短", "译文长度与原文显著不匹配"],
            ["日文残留", "译文中仍含日文字符"],
            ["数字/单位/符号遗漏", "原文数字未在译文出现"],
            ["术语遗漏", "词库词条未整词替换"],
            ["模型附带解释文字", "译文中混入模型说明性内容"]],
           col_widths=[2.6, 5.0], font_size=9)
card(s, 5.1, 5.62, 4.1, 1.2, "AI 润色", [
    "复用已有 DeepSeek 翻译服务",
    "自然/口语/热血/搞笑/正式/自定义",
    "失败时保留原译文",
], kind="accent", num="可选增强", title_size=11, body_size=9)
placeholder(s, 9.4, 5.62, 3.3, 1.24, "[图片占位-09-词库/质量告警界面-3:2]",
            "放入真实词库界面或翻译结果中质量提示的截图。")
notes(s, "翻译讲三件事：气泡整块一次翻译保留换行与句序；术语库整词替换保持一致；失败走回退链。"
         "右侧质量告警表逐项可指着说；AI 润色强调复用已有 DeepSeek 服务、失败保留原译文，不是新部署的模型。")

# ============================================================
# 第 10 页 · 译文如何落回画面：排版与人工编辑
# ============================================================
s = new_slide()
chrome(s, 10, "排版与编辑")
header(s, "RENDER & EDIT · 10", "译文如何落回画面：排版与人工编辑",
       "自动排版保证不越界，人工编辑保证可修正")
placeholder(s, MARGIN, 2.2, 5.9, 3.32, "[图片占位-10-结果编辑器与排版对照-16:9]",
            "放入真实大图预览、文本框选中、字体预览或擦除撤销界面。")
rect(s, MARGIN, 5.72, 5.9, 0.95, fill=INK)
tb(s, MARGIN + 0.16, 5.86, 5.6, 0.7, [
    {"text": "自动排版 + 人工微调闭环", "size": 11.5, "color": WHITE, "bold": True, "space_after": 3},
    {"text": "自动排版先给出不越界的译图，人工在编辑器里做最后修正", "size": 9, "color": PAPER_DIM},
])
card(s, 6.9, 2.2, 5.8, 1.28, "横排与竖排采用不同排版策略", [
    "横排：均衡断行 + 字号选择，规避孤字",
    "竖排：按气泡分列、按原文行序分配，控制列宽与字距",
], num="排版策略", title_size=11.5, body_size=9.5)
card(s, 6.9, 3.62, 5.8, 1.28, "渲染安全机制", [
    "文字与气泡掩膜覆盖率不足 50% 时回退矩形裁剪",
    "气泡几何不可靠时按安全扩展框、紧致框分级处理",
], num="防丢字", title_size=11.5, body_size=9.5)
card(s, 6.9, 5.04, 5.8, 1.63, "用户可操作编辑", [
    "拖拽移动文本框、调整尺寸",
    "修改字体 / 字号 / 颜色 / 方向并即时预览",
    "新增或删除文本框；局部擦除并支持撤销",
], num="人工兜底", title_size=11.5, body_size=9.5)
notes(s, "左边讲排版策略：横排均衡换行，竖排按气泡分列并保持原文顺序；覆盖率不足回退矩形裁剪。"
         "右边讲人工兜底：拖拽、改字体、局部擦除可撤销，形成自动加微调闭环，指向编辑器截图讲。")

# ============================================================
# 第 11 页 · 任务与批量：如何保证可用性
# ============================================================
s = new_slide()
chrome(s, 11, "任务与批量")
header(s, "BATCH · 11", "任务与批量：如何保证可用性",
       "批量是编排，推理仍串行；结果路径只信任服务端记录")
tb(s, MARGIN, 2.15, 8.0, 0.26, [{"text": "批量任务时序", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "spacing": 1.5}])
chevron_row(s, MARGIN, 2.45, CONTENT_W, 0.56,
            ["批量校验", "创建子任务", "逐张串行处理", "汇总平均进度", "生成 ZIP"],
            size=10, accent_last=True)
card(s, MARGIN, 3.35, 5.9, 1.5, "单 worker + 流水线锁", [
    "任务管理器默认单 worker，模型推理不并发",
    "批量只并发编排，不并发推理；ZIP 保留原文件名与顺序，重名自动改名，含 manifest.json 与 errors.txt",
], kind="ink", num="关键约束", title_size=11.5, body_size=9.5)
card(s, 6.85, 3.35, 5.85, 1.5, "安全边界", [
    "文件类型与大小校验；批量数量与总大小限制",
    "结果路径只从服务端任务记录解析，不接受客户端任意路径",
], num="输入与输出", title_size=11.5, body_size=9.5)
tb(s, MARGIN, 5.1, 6.0, 0.26, [{"text": "批量结果状态", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "spacing": 1.5}])
make_table(s, MARGIN, 5.4, 7.4, 1.42,
           ["状态", "表现", "用户操作"],
           [["正常", "全部子任务成功", "直接下载 ZIP"],
            ["部分成功", "失败子任务标明原因", "可单独重试后重下"],
            ["失败", "errors.txt 列出明细", "按提示修正后重来"]],
           col_widths=[1.5, 3.2, 2.7], font_size=9)
placeholder(s, 8.6, 5.15, 4.13, 1.75, "[图片占位-11-批量任务界面或下载清单-16:9]",
            "放入多图任务进度、部分失败提示或 ZIP 下载界面截图。")
notes(s, "强调批量是编排不是并发：校验、建子任务、逐张串行、汇总进度、生成 ZIP。安全边界要主动说："
         "结果路径只从服务端任务记录解析。下方状态表说明部分成功也可下载、可单独重试。")

# ============================================================
# 第 12 页 · 产品界面：从工具到任务工作台
# ============================================================
s = new_slide()
chrome(s, 12, "产品界面")
header(s, "PRODUCT UI · 12", "产品界面：从工具到任务工作台",
       "界面围绕“任务进行到哪、哪张失败、下一步做什么”组织")
ui_points = [
    "侧边栏：“翻译任务”与“专有名词库”同级",
    "主区域：接近对话式的任务流",
    "底部：上传 / 粘贴图片并提交",
    "每张图完成后显示预览、重试、下载入口",
]
uy = 2.25
for i, t in enumerate(ui_points):
    sp = rect(s, MARGIN, uy, 5.0, 0.72, fill=(GREY_1 if i % 2 == 0 else WHITE), line=(GREY_2 if i % 2 else None))
    shape_text(sp, [{"runs": [
        {"text": f"0{i + 1}  ", "size": 10, "color": ACCENT, "font": F_MONO, "bold": True},
        {"text": t, "size": 10.5, "color": INK, "bold": True},
    ], "align": PP_ALIGN.LEFT, "line_spacing": 1.15}])
    sp.text_frame.margin_left = Inches(0.14)
    uy += 0.86
tb(s, MARGIN, uy + 0.08, 5.0, 1.1, [
    {"text": "设计目的", "size": 10.5, "color": INK, "bold": True, "space_after": 3},
    {"text": "让用户知道当前任务进行到哪里、哪一张失败、下一步能做什么。", "size": 9.5, "color": GREY_3, "line_spacing": 1.3},
])
placeholder(s, 6.0, 2.25, 6.73, 3.79, "[图片占位-12-主界面/任务侧边栏/词库/编辑器拼版-16:9]",
            "放入实际前端页面截图（2—4 张拼版），统一裁剪比例和浏览器边框；截图旁加短标签。")
labels12 = ["任务状态", "图片预览", "文本编辑", "词库管理"]
for i, t in enumerate(labels12):
    sp = rect(s, 6.0 + i * 1.72, 6.2, 1.6, 0.5, fill=(ACCENT if i == 0 else GREY_1))
    shape_text(sp, [{"text": t, "size": 9.5, "color": (WHITE if i == 0 else INK), "bold": True}])
notes(s, "这页讲信息架构：侧边栏翻译任务与词库同级，主区域是对话式任务流。用截图拼版指认四个细节："
         "任务状态、图片预览、文本编辑、词库管理。设计目的：让用户始终知道下一步能做什么。")

# ============================================================
# 第 13 页 · AI 协作：人工决策与迭代证据
# ============================================================
s = new_slide()
chrome(s, 13, "AI 协作")
header(s, "AI COLLABORATION · 13", "AI 协作：人工决策与迭代证据",
       "人工做约束、核查与取舍；AI 做分析、草拟与局部实现")
card(s, MARGIN, 2.2, 3.2, 1.42, "人工负责", [
    "需求约束与优先级", "事实核查与运行验证", "体验判断与最终取舍",
], num="分工 A", title_size=11.5, body_size=9.5)
card(s, MARGIN + 3.4, 2.2, 3.2, 1.42, "AI 负责", [
    "代码阅读与问题定位", "方案草拟与局部实现", "测试建议",
], num="分工 B", title_size=11.5, body_size=9.5)
tb(s, MARGIN, 4.1, 6.6, 0.26, [{"text": "迭代循环", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "spacing": 1.5}])
chevron_row(s, MARGIN, 4.4, 6.6, 0.5, ["需求", "分析", "实现", "人工验证", "迭代"], size=9.5, accent_last=True)
card(s, MARGIN, 5.1, 6.6, 1.72, "AI 能力边界", [
    "不能替代真实运行；不能自动证明准确率",
    "不能把规划内容当成已实现；不能绕过人工对 UI 与输出效果的验收",
], kind="ink", num="边界声明", title_size=11.5, body_size=9.5)
ev = [
    ("E1", "先读代码与真实页面", "AI 先阅读代码和真实页面再提方案，避免凭空编造功能"),
    ("E2", "词库删除边界", "人工发现内置/用户词条边界，要求只删用户词条；AI 完成数据迁移与验证"),
    ("E3", "AI 润色约束", "人工限定复用已有 DeepSeek、不新增模型、不改主链，失败保留原译文"),
    ("E4", "可观察验收条件", "擦除、字体预览、文本框拖拽、文件夹上传：人工定验收条件并运行复核"),
]
for i, (nb, t, d) in enumerate(ev):
    x = 7.4 + (i % 2) * 2.72
    y = 2.2 + (i // 2) * 1.55
    card(s, x, y, 2.6, 1.42, t, [d], num=nb, title_size=10, body_size=8)
placeholder(s, 7.4, 5.35, 5.32, 1.47, "[图片占位-13-AI 协作原始记录截图-4:3]",
            "放入真实交互记录截图，优先选择有明确约束、问题发现和验证结果的记录。")
notes(s, "答辩关键页。分工先说清：人工定约束、查事实、跑验证、做取舍；AI 读代码、找问题、草拟方案、局部实现。"
         "四类证据选一两个展开讲，最后主动说 AI 能力边界四条，把验证主动权留在人手里。")

# ============================================================
# 第 14 页 · 怎么证明：测试与现场演示证据
# ============================================================
s = new_slide()
chrome(s, 14, "验证与演示")
header(s, "VERIFICATION · 14", "怎么证明：测试与现场演示证据",
       "只用真实验证结果说话：后端 171 项测试通过，前端生产构建成功")
sp = rect(s, MARGIN, 2.2, 3.6, 1.7, fill=GREY_1)
shape_text(sp, [
    {"text": "后端测试", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "space_after": 4, "spacing": 1.5},
    {"text": "171 项通过", "size": 24, "color": ACCENT, "font": F_TITLE, "space_after": 4},
    {"text": "stdlib unittest · 验证范围以当前代码和测试环境为准", "size": 8.5, "color": GREY_3, "line_spacing": 1.25},
])
sp = rect(s, MARGIN + 3.8, 2.2, 3.6, 1.7, fill=GREY_1)
shape_text(sp, [
    {"text": "前端生产构建", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "space_after": 4, "spacing": 1.5},
    {"text": "构建成功", "size": 24, "color": ACCENT, "font": F_TITLE, "space_after": 4},
    {"text": "npm run build · 验证范围以当前代码和测试环境为准", "size": 8.5, "color": GREY_3, "line_spacing": 1.25},
])
card(s, 8.6, 2.2, 4.13, 1.7, "未提供的指标", [
    "无统一数据集测量的 OCR 准确率",
    "无统一数据集测量的翻译准确率",
    "无平均耗时统计，不做虚构比较",
], kind="outline", num="诚实声明", title_size=11, body_size=9)
tb(s, MARGIN, 4.15, 8.0, 0.26, [{"text": "现场演示顺序（时间不足则播放录屏）", "size": 10, "color": GREY_3, "font": F_MONO, "bold": True, "spacing": 1.5}])
demo14 = ["选日文/英文样例", "上传/粘贴", "查看进度", "展示译图", "打开编辑", "修改一处", "下载结果"]
n14 = len(demo14)
gap14 = 0.07
cw14 = (CONTENT_W - gap14 * (n14 - 1)) / n14
for i, t in enumerate(demo14):
    x = MARGIN + i * (cw14 + gap14)
    sp = rect(s, x, 4.45, cw14, 0.56, fill=(GREY_1 if i < n14 - 1 else ACCENT),
              shape=(MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON))
    shape_text(sp, [{"text": f"{i + 1:02d}  {t}", "size": 9, "color": (INK if i < n14 - 1 else WHITE), "bold": True, "line_spacing": 1.1}])
placeholder(s, MARGIN, 5.3, 6.6, 1.6, "[图片占位-14-测试结果与真实译图对照-16:9]",
            "左侧放测试通过截图，右侧放仓库真实日文/英文样例的原译对照。")
tb(s, 7.5, 5.3, 5.2, 1.6, [
    {"text": "三类证据位", "size": 10.5, "color": INK, "bold": True, "space_after": 4},
    {"text": "· 测试命令结果截图", "size": 9.5, "color": GREY_3, "space_after": 2},
    {"text": "· 原图 / 译图对照", "size": 9.5, "color": GREY_3, "space_after": 2},
    {"text": "· 实际操作录屏或 UI 截图", "size": 9.5, "color": GREY_3},
])
notes(s, "只讲真实证据：后端 171 项测试通过、前端构建成功，并注明验证范围以当前代码和环境为准。"
         "主动声明没有统一数据集下的准确率与耗时，不做虚构比较。然后按七步演示顺序走一遍，时间不足就播放录屏。")

# ============================================================
# 第 15 页 · 结论、边界与答辩入口（统一浅色风格）
# ============================================================
s = new_slide()
chrome(s, 15, "结论与答辩")
header(s, "CONCLUSION · 15", "结论、边界与答辩入口",
       "闭环已通、难点已解、边界已明")
summary15 = [
    ("01", "已实现“图片输入—识别—修复—翻译—排版—编辑—下载”的完整闭环"),
    ("02", "核心难点在漫画场景下的方向识别、气泡边界、背景修复和可控回退"),
    ("03", "下一步优先完善历史记录、章节级管理和更多真实样例验证"),
]
sy = 2.2
for nb, t in summary15:
    tb(s, MARGIN, sy, 8.0, 0.62, [
        {"runs": [
            {"text": nb + "  ", "size": 13, "color": ACCENT, "font": F_MONO, "bold": True},
            {"text": t, "size": 12.5, "color": INK},
        ], "line_spacing": 1.2},
    ])
    sy += 0.68
box = rect(s, MARGIN, 4.35, 4.6, 1.95, fill=GREY_1, line=GREY_2, line_w=0.75)
shape_text(box, [
    {"text": "当前可验收能力", "size": 11.5, "color": INK, "bold": True, "align": PP_ALIGN.LEFT, "space_after": 5},
    *[{"text": "· " + t, "size": 9.5, "color": GREY_3, "align": PP_ALIGN.LEFT, "space_after": 3, "line_spacing": 1.2}
      for t in ["图片/文件夹输入与日英→中翻译", "任务进度、预览、单张重试", "词库、文本框编辑、局部擦除撤销", "单张/批量 ZIP 下载"]],
], anchor=MSO_ANCHOR.TOP)
box.text_frame.margin_left = Inches(0.16)
box.text_frame.margin_top = Inches(0.12)
box = rect(s, MARGIN + 4.8, 4.35, 3.2, 1.95, fill=WHITE, line=GREY_2, line_w=0.75)
shape_text(box, [
    {"text": "后续可扩展方向", "size": 11.5, "color": INK, "bold": True, "align": PP_ALIGN.LEFT, "space_after": 5},
    *[{"text": "· " + t, "size": 9.5, "color": GREY_3, "align": PP_ALIGN.LEFT, "space_after": 3, "line_spacing": 1.2}
      for t in ["历史翻译记录", "章节级管理", "EPUB / CBZ / PDF", "团队协作、云端 GPU"]],
], anchor=MSO_ANCHOR.TOP)
box.text_frame.margin_left = Inches(0.16)
box.text_frame.margin_top = Inches(0.12)
tb(s, MARGIN, 6.45, 8.0, 0.5, [
    {"text": "谢谢，欢迎提问。", "size": 18, "color": INK, "font": F_TITLE},
])
rect(s, 9.1, 2.2, 3.63, 0.001, fill=None)
placeholder(s, 9.6, 2.2, 2.6, 2.6, "[图片占位-15-最终译图或项目二维码-1:1]",
            "可放一张最有代表性的真实译图；如有项目地址再替换为真实二维码。")
sp = rect(s, 9.6, 5.1, 2.6, 1.2, fill=ACCENT)
shape_text(sp, [
    {"text": "现场演示入口", "size": 12, "color": WHITE, "bold": True, "space_after": 3},
    {"text": "打开系统实操或播放录屏", "size": 9, "color": RGBColor(0xD5, 0xDD, 0xF5)},
])
notes(s, "三句话收束：闭环已通、难点在哪、下一步做什么。两栏对照可验收能力与扩展方向。致谢后进入提问，"
         "若老师要看现场，直接打开系统演示或播放录屏。\n"
         "素材替换清单：01 首页或原译对照｜02 复杂文字局部｜04 上传/进度/预览截图｜05 架构图或目录截图｜"
         "06 流程或日志截图｜07 OCR 检测框｜08 四格对照｜09 词库/质量告警｜10 编辑器界面｜11 批量任务界面｜"
         "12 界面拼版｜13 AI 协作记录｜14 测试结果+原译对照｜15 代表性译图或二维码。")

prs.save(OUT_FILE)
print(f"OK -> {OUT_FILE}")
