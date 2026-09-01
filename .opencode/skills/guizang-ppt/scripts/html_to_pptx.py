#!/usr/bin/env python3
"""
html_to_pptx.py - 将 HTML 幻灯片转换为可编辑的 PowerPoint (.pptx)

参考 huashu-slides 的 html2pptx.js 设计思路：
1. 递归遍历所有语义元素（P, H1-H6, UL, OL, LI, DIV）
2. 为每个语义元素创建独立的 PPTX 文本框
3. 支持内联格式（<b>, <i>, <span> 等）
4. 完整的样式提取（fontSize, color, bold, italic 等）
5. CSS 布局解析（grid, flex, span 类名）

依赖：python-pptx (pip install python-pptx)
"""

import re
import sys
from pathlib import Path
from typing import Optional, List, Tuple, Dict, Any
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup, Tag, NavigableString


# ============================================================
# 工具函数
# ============================================================

def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """HEX 颜色转 RGB"""
    hex_color = hex_color.strip('#')
    if len(hex_color) == 6:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    elif len(hex_color) == 3:
        return tuple(int(c*2, 16) for c in hex_color)
    return (0, 0, 0)


def parse_css_color(color_str: str) -> Optional[RGBColor]:
    """解析 CSS 颜色值"""
    if not color_str:
        return None
    color_str = color_str.strip().lower()

    # HEX 颜色
    if color_str.startswith('#'):
        try:
            r, g, b = hex_to_rgb(color_str)
            return RGBColor(r, g, b)
        except:
            pass

    # RGB/RGBA
    rgb_match = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', color_str)
    if rgb_match:
        return RGBColor(int(rgb_match.group(1)),
                       int(rgb_match.group(2)),
                       int(rgb_match.group(3)))

    # 颜色名称映射
    color_map = {
        'white': (255, 255, 255),
        'black': (0, 0, 0),
        'red': (255, 0, 0),
        'green': (0, 128, 0),
        'blue': (0, 0, 255),
        'yellow': (255, 255, 0),
        'gray': (128, 128, 128),
        'grey': (128, 128, 128),
    }
    if color_str in color_map:
        r, g, b = color_map[color_str]
        return RGBColor(r, g, b)

    return None


def parse_font_size(size_str: str) -> int:
    """解析字体大小，默认 18pt"""
    if not size_str:
        return 18
    size_str = size_str.strip().lower()

    # 移除 px/em/pt 单位
    size_str = re.sub(r'[a-z%]+', '', size_str)
    try:
        size = float(size_str)
        # px 转 pt（假设 1px = 0.75pt）
        if 'px' in size_str.lower() or not size_str:
            size = size * 0.75
        return max(8, min(200, int(size)))
    except:
        return 18


def parse_inline_style(style_str: str) -> dict:
    """解析 inline style 属性"""
    styles = {}
    if not style_str:
        return styles

    for item in style_str.split(';'):
        if ':' in item:
            prop, value = item.split(':', 1)
            styles[prop.strip().lower()] = value.strip()

    return styles


# ============================================================
# CSS 布局解析
# ============================================================

@dataclass
class LayoutRect:
    """元素布局矩形"""
    left: float = 0.0   # 英寸
    top: float = 0.0
    width: float = 13.333  # 16:9 幻灯片宽度
    height: float = 7.5    # 幻灯片高度


@dataclass
class ParsedElement:
    """解析后的元素（包含位置和样式）"""
    tag: str
    content: Any  # str 或 List[dict] 用于带格式的文本
    rect: LayoutRect = field(default_factory=LayoutRect)
    style: Dict[str, Any] = field(default_factory=dict)
    element_type: str = 'text'  # text, list, image, shape


class LayoutCalculator:
    """CSS 布局计算器"""
    SLIDE_WIDTH = 13.333
    SLIDE_HEIGHT = 7.5
    MARGIN_LEFT = 0.8
    MARGIN_RIGHT = 0.8
    MARGIN_TOP = 0.6
    CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

    # 风格 A 的类名 → 布局模式
    STYLE_A_LAYOUT_PATTERNS = {
        'grid-2-7-5': [(2, 0.17), (7, 0.33), (5, 0.58)],
        'grid-2-6-6': [(2, 0.17), (6, 0.33), (6, 0.50)],
        'grid-2-8-4': [(2, 0.17), (8, 0.33), (4, 0.67)],
        'grid-3-3': [(3, 0.0), (3, 0.25), (3, 0.50)],
        'grid-6': [(6, 0.0)],
        'grid-3': [(3, 0.0)],
        'grid-4': [(4, 0.0)],
    }

    # 风格 B 的 span 类名映射
    SPAN_TO_COLS = {
        'span-1': 1, 'span-2': 2, 'span-3': 3, 'span-4': 4,
        'span-5': 5, 'span-6': 6, 'span-7': 7, 'span-8': 8,
        'span-9': 9, 'span-10': 10, 'span-11': 11, 'span-12': 12,
    }

    @classmethod
    def detect_layout_pattern(cls, classes: List[str]) -> Optional[str]:
        """检测布局模式"""
        for cls_name in classes:
            if cls_name.startswith('grid-'):
                return cls_name
        return None

    @classmethod
    def get_position_from_layout(cls, layout_pattern: str, index: int,
                                  total_items: int, style: str = "A") -> LayoutRect:
        """根据布局模式计算元素位置"""
        rect = LayoutRect()

        if style == "A":
            pattern = cls.STYLE_A_LAYOUT_PATTERNS.get(layout_pattern)
            if pattern and index < len(pattern):
                cols, start_ratio = pattern[index]
                rect.left = cls.MARGIN_LEFT + cls.CONTENT_WIDTH * start_ratio
                rect.width = cls.CONTENT_WIDTH * (cols / 12) - 0.1
                rect.top = cls.MARGIN_TOP
                rect.height = cls.SLIDE_HEIGHT - cls.MARGIN_TOP - 0.5
            else:
                rect.left = cls.MARGIN_LEFT
                rect.width = cls.CONTENT_WIDTH
                rect.top = cls.MARGIN_TOP + (index * 1.2)
                rect.height = 1.0
        else:
            rect.left = cls.MARGIN_LEFT
            rect.width = cls.CONTENT_WIDTH
            rect.top = cls.MARGIN_TOP + (index * 1.0)
            rect.height = 0.8

        return rect


# ============================================================
# 递归 HTML 解析器
# ============================================================

class RecursiveHTMLElementParser:
    """
    递归 HTML 解析器

    参考 huashu-slides 的 html2pptx.js：
    1. 递归遍历所有语义元素
    2. 为每个元素提取文本和样式
    3. 支持内联格式
    """

    # 文本元素标签
    TEXT_TAGS = {'p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'span', 'div'}

    # 列表标签
    LIST_TAGS = {'ul', 'ol'}

    def __init__(self, html_content: str, style: str = "A"):
        self.html = html_content
        self.style = style
        self.soup = BeautifulSoup(html_content, 'html.parser')

    def get_classes(self, element: Tag) -> List[str]:
        """获取元素的类名列表"""
        return element.get('class', [])

    def get_style(self, element: Tag) -> Dict[str, str]:
        """获取元素的 inline style"""
        style_str = element.get('style', '')
        return parse_inline_style(style_str)

    def is_hidden(self, element: Tag) -> bool:
        """检查元素是否隐藏"""
        if element.get('data-html2pptx-ignore'):
            return True
        style = self.get_style(element)
        if style.get('display') == 'none':
            return True
        if style.get('visibility') == 'hidden':
            return True
        if style.get('opacity') == '0':
            return True
        return False

    def extract_text_with_formatting(self, element: Tag) -> Any:
        """
        递归提取文本及其格式

        返回：
        - str: 纯文本
        - List[dict]: 带格式的文本段列表 [{"text": "...", "options": {...}}, ...]
        """
        runs = []

        def process_node(node, base_options=None):
            if base_options is None:
                base_options = {}

            # 处理文本节点
            if isinstance(node, NavigableString):
                text = str(node).strip()
                if text:
                    # 合并相邻的纯文本
                    if runs and 'options' not in runs[-1]:
                        runs[-1] = runs[-1] + text
                    else:
                        runs.append(text)
                return

            # 处理元素节点
            if not isinstance(node, Tag):
                return

            tag_name = node.name.lower()

            # 获取计算样式（从 style 属性简单推断）
            style = self.get_style(node)
            options = {**base_options}

            # 处理内联格式标签
            if tag_name in ('b', 'strong', 'span'):
                if style.get('font-weight') in ('bold', '600', '700') or 'bold' in str(style):
                    options['bold'] = True
            if tag_name in ('i', 'em'):
                options['italic'] = True
            if tag_name in ('u', ):
                options['underline'] = True

            # 处理颜色
            if 'color' in style:
                color = parse_css_color(style['color'])
                if color:
                    options['color'] = color

            # 处理字体大小
            if 'font-size' in style:
                size = parse_font_size(style['font-size'])
                if size:
                    options['fontSize'] = size

            # 处理换行
            if tag_name == 'br':
                if runs:
                    runs.append('\n')
                return

            # 递归处理子节点
            for child in node.children:
                process_node(child, options)

        process_node(element)

        # 清理并合并
        if not runs:
            return ''

        # 如果全是纯文本，返回字符串
        if all(isinstance(r, str) for r in runs):
            return ' '.join(runs)

        # 否则返回带格式的段列表
        result = []
        for run in runs:
            if isinstance(run, str):
                if result and 'options' not in result[-1]:
                    result[-1] = result[-1] + ' ' + run
                else:
                    result.append(run)
            else:
                result.append(run)
        return result

    def extract_list_items(self, element: Tag) -> List[Any]:
        """提取列表项"""
        items = []
        li_elements = element.find_all('li', recursive=False)

        if not li_elements:
            # 查找嵌套的 li
            li_elements = element.find_all('li')

        for i, li in enumerate(li_elements):
            text = self.extract_text_with_formatting(li)
            is_last = (i == len(li_elements) - 1)

            item = {
                'text': text if isinstance(text, str) else text,
                'options': {
                    'breakLine': not is_last,
                    'bullet': True
                }
            }
            items.append(item)

        return items

    def parse_element(self, element: Tag, index: int = 0) -> Optional[ParsedElement]:
        """解析单个元素"""

        if self.is_hidden(element):
            return None

        tag_name = element.name.lower()
        classes = self.get_classes(element)
        style = self.get_style(element)

        # 检测是否有背景或边框（作为 shape 处理）
        has_bg = any(k for k in style.keys() if 'background' in k or 'background-color' in k)
        border_keys = ['border', 'border-width', 'border-top-width', 'border-right-width',
                       'border-bottom-width', 'border-left-width']
        has_border = any(style.get(k) for k in border_keys)

        # 根据标签类型处理
        if tag_name in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
            # 标题元素
            text = self.extract_text_with_formatting(element)
            if not text:
                return None

            # 计算字体大小
            font_sizes = {
                'h1': 44, 'h2': 36, 'h3': 28, 'h4': 24, 'h5': 20, 'h6': 18
            }
            font_size = font_sizes.get(tag_name, 24)

            # 检测是否加粗
            is_bold = True  # 标题默认加粗

            return ParsedElement(
                tag=tag_name,
                content=text,
                style={
                    'fontSize': font_size,
                    'bold': is_bold,
                    'italic': False,
                    'underline': False,
                    'align': 'left'
                },
                element_type='text'
            )

        elif tag_name == 'p':
            # 段落元素
            text = self.extract_text_with_formatting(element)
            if not text:
                return None

            return ParsedElement(
                tag=tag_name,
                content=text,
                style={
                    'fontSize': 16,
                    'bold': False,
                    'italic': False,
                    'underline': False,
                    'align': 'left'
                },
                element_type='text'
            )

        elif tag_name in ('ul', 'ol'):
            # 列表
            items = self.extract_list_items(element)
            if not items:
                return None

            return ParsedElement(
                tag=tag_name,
                content=items,
                style={
                    'fontSize': 14,
                    'align': 'left'
                },
                element_type='list'
            )

        elif tag_name == 'img':
            # 图片
            src = element.get('src', '')
            if not src:
                return None

            return ParsedElement(
                tag='img',
                content=src,
                element_type='image'
            )

        elif tag_name == 'div' and (has_bg or has_border):
            # 带样式的容器作为 shape
            return ParsedElement(
                tag='div',
                content='',
                style={
                    'fill': style.get('background-color'),
                    'border': has_border
                },
                element_type='shape'
            )

        return None

    def parse_slide(self, slide: Tag) -> List[ParsedElement]:
        """解析幻灯片的所有元素"""
        elements = []

        # 递归查找所有语义元素
        for tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'img', 'div']:
            for element in slide.find_all(tag_name, recursive=False):
                elem = self.parse_element(element, len(elements))
                if elem:
                    elements.append(elem)

        return elements


# ============================================================
# 主题配色配置
# ============================================================

STYLE_A_THEMES = {
    "墨水经典": {
        "primary": RGBColor(0x1A, 0x1A, 0x1A),
        "paper": RGBColor(0xF8, 0xF5, 0xF0),
        "accent": RGBColor(0x4A, 0x4A, 0x4A),
        "text": RGBColor(0x33, 0x33, 0x33),
    },
    "靛蓝瓷": {
        "primary": RGBColor(0x1A, 0x36, 0x5C),
        "paper": RGBColor(0xF0, 0xF4, 0xF8),
        "accent": RGBColor(0x4A, 0x90, 0xD9),
        "text": RGBColor(0x2C, 0x3E, 0x50),
    },
    "森林墨": {
        "primary": RGBColor(0x1A, 0x3A, 0x2C),
        "paper": RGBColor(0xF0, 0xF5, 0xF0),
        "accent": RGBColor(0x4A, 0x7C, 0x59),
        "text": RGBColor(0x2C, 0x3E, 0x35),
    },
    "牛皮纸": {
        "primary": RGBColor(0x5C, 0x40, 0x33),
        "paper": RGBColor(0xF5, 0xF0, 0xE5),
        "accent": RGBColor(0x8B, 0x69, 0x4F),
        "text": RGBColor(0x3C, 0x2E, 0x22),
    },
    "沙丘": {
        "primary": RGBColor(0x4A, 0x4A, 0x45),
        "paper": RGBColor(0xFA, 0xF8, 0xF5),
        "accent": RGBColor(0xC9, 0xA8, 0x6C),
        "text": RGBColor(0x3C, 0x3C, 0x38),
    },
}

STYLE_B_THEMES = {
    "IKB蓝": {"accent": RGBColor(0x00, 0x28, 0xFF)},
    "柠檬黄": {"accent": RGBColor(0xFF, 0xED, 0x00)},
    "柠檬绿": {"accent": RGBColor(0x84, 0xC1, 0x00)},
    "安全橙": {"accent": RGBColor(0xFF, 0x6B, 0x00)},
}


# ============================================================
# PPT 生成函数
# ============================================================

def create_prs(style: str = "A") -> Presentation:
    """创建 PPT 对象"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs: Presentation) -> object:
    """添加空白幻灯片"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_slide_background(slide, color: RGBColor):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text: Any, left: float, top: float,
                 width: float, height: float,
                 font_size: int = 18, font_name: str = "Microsoft YaHei",
                 color: RGBColor = None, bold: bool = False,
                 italic: bool = False,
                 alignment: PP_ALIGN = PP_ALIGN.LEFT) -> object:
    """添加文本框"""
    if not text:
        return None

    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    # 处理多行文本
    if isinstance(text, str):
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(font_size)
            p.font.name = font_name
            if color:
                p.font.color.rgb = color
            p.font.bold = bold
            p.font.italic = italic
            p.alignment = alignment
    else:
        # 带格式的文本段列表
        for i, run in enumerate(text):
            if isinstance(run, str):
                run = {'text': run}

            text_content = run.get('text', '')
            if not text_content:
                continue

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text_content
            p.font.size = Pt(run.get('fontSize', font_size))
            p.font.name = font_name
            if color:
                p.font.color.rgb = color
            p.font.bold = run.get('bold', bold)
            p.font.italic = run.get('italic', italic)
            p.alignment = run.get('align', alignment)

    return txBox


def add_shape(slide, shape_type, left: float, top: float,
              width: float, height: float, fill_color: RGBColor = None,
              line_color: RGBColor = None, line_width: float = 0):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color and line_width > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    return shape


def add_list(slide, items: List, left: float, top: float,
             width: float, height: float,
             font_size: int = 14, font_name: str = "Microsoft YaHei",
             color: RGBColor = None):
    """添加列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if isinstance(item, dict):
            text = item.get('text', '')
            options = item.get('options', {})
        else:
            text = str(item)
            options = {}

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {text}" if isinstance(text, str) else str(text)
        p.font.size = Pt(font_size)
        p.font.name = font_name
        if color:
            p.font.color.rgb = color

    return txBox


def add_image(slide, image_path: str, left: float, top: float,
              width: float, height: float):
    """添加图片"""
    try:
        slide.shapes.add_picture(image_path,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    except Exception as e:
        print(f"警告: 图片加载失败 {image_path}: {e}")


# ============================================================
# 幻灯片渲染
# ============================================================

def render_slide(prs: Presentation, slide_data: dict, theme_config: dict, style: str = "A"):
    """渲染一张幻灯片"""
    theme = slide_data["theme"]
    elements = slide_data["elements"]

    slide = add_blank_slide(prs)

    # 根据主题设置背景
    if theme in ["hero dark", "dark"]:
        set_slide_background(slide, theme_config["primary"])
        text_color = RGBColor(0xFF, 0xFF, 0xFF)
        accent_color = theme_config["accent"]
    else:
        set_slide_background(slide, theme_config["paper"])
        text_color = theme_config["text"]
        accent_color = theme_config["primary"]

    if theme.startswith("hero"):
        # Hero 页面：大标题居中
        main_title = None
        for elem in elements:
            if elem.tag in ('h1', 'h2', 'h3'):
                main_title = elem.content
                break

        if main_title:
            title_text = main_title if isinstance(main_title, str) else \
                        (main_title[0].get('text', '') if main_title else '')
            add_text_box(slide, title_text,
                        left=1, top=2.5, width=11.33, height=1.5,
                        font_size=44, font_name="Microsoft YaHei",
                        color=text_color, bold=True,
                        alignment=PP_ALIGN.CENTER)
    else:
        # 普通页面：渲染所有元素
        layout_calc = LayoutCalculator()

        for i, elem in enumerate(elements):
            # 使用布局计算确定位置
            layout_pattern = None
            if elem.tag in ('div', 'section'):
                classes = elem.style.get('classes', [])
                layout_pattern = layout_calc.detect_layout_pattern(classes)

            if layout_pattern:
                rect = layout_calc.get_position_from_layout(layout_pattern, i, len(elements), style)
            else:
                # 默认位置：垂直堆叠
                rect = LayoutRect()
                rect.left = layout_calc.MARGIN_LEFT
                rect.width = layout_calc.CONTENT_WIDTH
                rect.top = layout_calc.MARGIN_TOP + (i * 1.0)
                rect.height = 0.8

            # 渲染元素
            if elem.element_type == 'text':
                content = elem.content
                elem_style = elem.style

                font_size = elem_style.get('fontSize', 16)
                bold = elem_style.get('bold', False)
                italic = elem_style.get('italic', False)
                align = elem_style.get('align', 'left')

                # 对齐方式
                if align == 'center':
                    alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    alignment = PP_ALIGN.RIGHT
                else:
                    alignment = PP_ALIGN.LEFT

                add_text_box(slide, content,
                           left=rect.left, top=rect.top,
                           width=rect.width, height=rect.height,
                           font_size=font_size, font_name="Microsoft YaHei",
                           color=text_color, bold=bold, italic=italic,
                           alignment=alignment)

            elif elem.element_type == 'list':
                add_list(slide, elem.content,
                        left=rect.left, top=rect.top,
                        width=rect.width, height=rect.height,
                        font_size=elem.style.get('fontSize', 14),
                        color=text_color)

            elif elem.element_type == 'image':
                add_image(slide, elem.content,
                         left=rect.left, top=rect.top,
                         width=rect.width, height=rect.height)

            elif elem.element_type == 'shape':
                fill_color = None
                if elem.style.get('fill'):
                    fill_color = parse_css_color(elem.style['fill'])

                add_shape(slide, MSO_SHAPE.RECTANGLE,
                         left=rect.left, top=rect.top,
                         width=rect.width, height=rect.height,
                         fill_color=fill_color)


# ============================================================
# HTML 解析
# ============================================================

def extract_slides(html_content: str) -> List[dict]:
    """从 HTML 中提取所有幻灯片"""
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = []

    # 查找所有 section.slide
    for slide in soup.find_all('section', class_='slide'):
        slide_dict = {
            "theme": "light",
            "elements": []
        }

        # 检测主题
        class_list = slide.get('class', [])
        if 'hero' in class_list and 'dark' in class_list:
            slide_dict["theme"] = "hero dark"
        elif 'hero' in class_list:
            slide_dict["theme"] = "hero light"
        elif 'dark' in class_list:
            slide_dict["theme"] = "dark"

        # 检测布局
        layout = slide.get('data-layout', '')
        slide_dict["layout"] = layout

        # 解析元素
        parser = RecursiveHTMLElementParser(str(slide))

        for tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'img', 'div']:
            for element in slide.find_all(tag_name, recursive=False):
                elem = parser.parse_element(element)
                if elem:
                    slide_dict["elements"].append(elem)

        slides.append(slide_dict)

    return slides


# ============================================================
# 主转换函数
# ============================================================

def html_to_pptx(input_html: str, output_pptx: str,
                 style: str = "A", theme_name: str = "墨水经典") -> bool:
    """
    将 HTML 文件转换为 PPTX

    Args:
        input_html: HTML 文件路径或 HTML 内容字符串
        output_pptx: 输出 PPTX 文件路径
        style: "A" 电子杂志风 或 "B" 瑞士国际主义风
        theme_name: 主题名称

    Returns:
        True if successful, False otherwise
    """
    try:
        # 读取 HTML
        if Path(input_html).exists():
            html_content = Path(input_html).read_text(encoding='utf-8')
        else:
            html_content = input_html

        # 获取主题配色
        if style == "A":
            theme_config = STYLE_A_THEMES.get(theme_name, STYLE_A_THEMES["墨水经典"])
        else:
            theme_config = {"accent": STYLE_B_THEMES.get(theme_name, STYLE_B_THEMES["IKB蓝"])["accent"]}
            theme_config["primary"] = RGBColor(0x1A, 0x1A, 0x1A)
            theme_config["paper"] = RGBColor(0xFF, 0xFF, 0xFF)
            theme_config["text"] = RGBColor(0x1A, 0x1A, 0x1A)

        # 提取幻灯片
        slides_data = extract_slides(html_content)

        if not slides_data:
            print("警告: 未找到幻灯片内容")
            return False

        # 创建 PPT
        prs = create_prs(style)

        # 渲染每张幻灯片
        for slide_data in slides_data:
            render_slide(prs, slide_data, theme_config, style)

        # 保存
        prs.save(output_pptx)
        print(f"✓ 已生成: {output_pptx}")
        return True

    except Exception as e:
        import traceback
        print(f"错误: {e}")
        traceback.print_exc()
        return False


# ============================================================
# 命令行接口
# ============================================================

def main():
    if len(sys.argv) < 3:
        print("用法: python html_to_pptx.py <输入.html> <输出.pptx> [风格(A|B)] [主题]")
        print("\n风格 A 主题: 墨水经典, 靛蓝瓷, 森林墨, 牛皮纸, 沙丘")
        print("风格 B 主题: IKB蓝, 柠檬黄, 柠檬绿, 安全橙")
        print("\n特性:")
        print("  - 递归遍历所有语义元素（P, H1-H6, UL, OL, LI）")
        print("  - 为每个元素创建独立的 PPTX 文本框")
        print("  - 支持内联格式（<b>, <i>, <span> 等）")
        print("  - CSS Grid 布局解析")
        print("  - 支持 <img> 图片转换")
        print("  - 支持带背景/边框的 <div> 形状转换")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]
    style = sys.argv[3].upper() if len(sys.argv) > 3 else "A"
    theme = sys.argv[4] if len(sys.argv) > 4 else ("IKB蓝" if style == "B" else "墨水经典")

    if style not in ('A', 'B'):
        print("错误: 风格必须是 A 或 B")
        sys.exit(1)

    success = html_to_pptx(input_file, output_file, style, theme)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
