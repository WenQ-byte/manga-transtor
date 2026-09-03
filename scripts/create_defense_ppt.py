from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "AI_Manga_Translator_软件工程项目答辩.pptx"
FONT = "Microsoft YaHei"
MONO = "Consolas"
W, H = 13.333, 7.5
BG = RGBColor(10, 15, 27)
PANEL = RGBColor(20, 29, 48)
PANEL2 = RGBColor(28, 39, 63)
WHITE = RGBColor(240, 245, 252)
MUTED = RGBColor(157, 173, 199)
CYAN = RGBColor(40, 220, 240)
PURPLE = RGBColor(176, 111, 255)
PINK = RGBColor(245, 103, 198)
GREEN = RGBColor(80, 220, 170)
ORANGE = RGBColor(255, 178, 89)

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

def rgb(hexv):
    hexv = hexv.lstrip('#')
    return RGBColor.from_string(hexv.upper())

def shape(slide, typ, x, y, w, h, fill=None, line=None, radius=False):
    s = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    return s

def text(slide, value, x, y, w, h, size=18, color=WHITE, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.04):
    if not isinstance(font, str):
        align, font = font, FONT
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def add_bg(slide, title, kicker=None, page=None):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BG
    # understated grid / glow points
    for x in [0.6, 3.6, 6.6, 9.6, 12.6]:
        shape(slide, MSO_SHAPE.RECTANGLE, x, 0.6, 0.012, 6.3, line=None, fill=rgb("18233A"))
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.06, fill=CYAN)
    text(slide, "AI MANGA TRANSLATOR  /  SOFTWARE ENGINEERING DEFENSE", 0.6, 0.26, 8, 0.22, 8, MUTED, False, MONO)
    if kicker: text(slide, kicker.upper(), 0.62, 0.95, 4, 0.3, 10, CYAN, True, MONO)
    text(slide, title, 0.6, 1.22, 11.8, 0.56, 27, WHITE, True)
    if page is not None: text(slide, f"{page:02d}", 12.15, 0.27, 0.55, 0.25, 10, MUTED, False, MONO, PP_ALIGN.RIGHT)

def pill(slide, label, x, y, w, color=CYAN):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.34, fill=rgb("16283C"), line=color)
    text(slide, label, x, y+0.01, w, 0.3, 9, color, True, MONO, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body, accent=CYAN, body_size=12):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=PANEL, line=rgb("2A3A59"))
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, fill=accent)
    text(slide, title, x+0.22, y+0.17, w-0.4, 0.3, 14, WHITE, True)
    text(slide, body, x+0.22, y+0.58, w-0.4, h-0.72, body_size, MUTED, False, FONT, valign=MSO_ANCHOR.TOP)

def fit_picture(slide, path, x, y, w, h, pad=0.0):
    path = str(path)
    with Image.open(path) as im: iw, ih = im.size
    box_ratio = w / h; im_ratio = iw / ih
    if im_ratio > box_ratio:
        crop_h = ih; crop_w = int(ih * box_ratio); left = (iw-crop_w)//2; top = 0
    else:
        crop_w = iw; crop_h = int(iw / box_ratio); left = 0; top = (ih-crop_h)//2
    pic = slide.shapes.add_picture(path, Inches(x+pad), Inches(y+pad), Inches(w-2*pad), Inches(h-2*pad))
    pic.crop_left = left/iw; pic.crop_right = (iw-left-crop_w)/iw
    pic.crop_top = top/ih; pic.crop_bottom = (ih-top-crop_h)/ih
    return pic

def arrow(slide, x1, y1, x2, y2, color=CYAN, width=2.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color; line.line.width = Pt(width); line.line.end_arrowhead = True
    return line

def icon_circle(slide, label, x, y, color=CYAN, size=0.72):
    shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill=rgb("172840"), line=color)
    text(slide, label, x, y+0.01, size, size-0.02, 15, color, True, MONO, PP_ALIGN.CENTER)

def note(slide, value):
    text(slide, value, 0.62, 6.95, 12, 0.24, 8, MUTED)

def image_frame(slide, path, x, y, w, h, label):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=PANEL, line=rgb("304262"))
    fit_picture(slide, path, x+0.08, y+0.08, w-0.16, h-0.16)
    pill(slide, label, x+0.18, y+0.16, max(1.0, len(label)*0.12+0.38), PINK)

JP = ROOT / "test_image_japanese" / "e19eca8c467765ba30042abd639840e6.jpg"
EN = ROOT / "test_image_english" / "屏幕截图 2026-08-31 114913.png"
ZH = ROOT / "test_image_chinese" / "屏幕截图 2026-08-31 145655.png"
LOGO = ROOT / "frontend" / "public" / "comicx-app-icon.png"

# 1 cover
s = prs.slides.add_slide(blank); add_bg(s, "AI Manga Translator", "software engineering project defense", 1)
text(s, "基于 OCR、AI 翻译与图像修复的智能漫画翻译系统", 0.65, 2.05, 7.5, 0.55, 22, WHITE, True)
text(s, "把“识别、翻译、擦除、修复、排版”串成可运行的工程化流水线", 0.68, 2.75, 6.7, 0.38, 13, MUTED)
pill(s, "React / FastAPI / MIT OCR / OpenCV / LaMa", 0.68, 3.45, 4.0, CYAN)
fit_picture(s, LOGO, 9.5, 1.45, 2.65, 2.65)
shape(s, MSO_SHAPE.RECTANGLE, 8.6, 4.65, 3.8, 0.02, fill=rgb("2A3A59"))
text(s, "软件工程项目答辩", 8.62, 4.9, 3.5, 0.4, 15, CYAN, True)
text(s, "姓名 / 小组：________________\n日期：2026", 8.62, 5.38, 3.4, 0.68, 12, MUTED)

# 2 problem
s = prs.slides.add_slide(blank); add_bg(s, "为什么需要漫画翻译工程化？", "01 / problem", 2)
text(s, "传统流程把多个专业工具串在一起，效率与一致性都被人工操作限制。", 0.66, 1.95, 11.2, 0.35, 16, WHITE)
steps = [("识别文字", "OCR / 人工抄录", ORANGE), ("人工翻译", "逐句处理", PINK), ("擦除原文", "Photoshop", PURPLE), ("修复背景", "局部重建", CYAN), ("重新排版", "逐框调整", GREEN)]
for i,(a,b,c) in enumerate(steps):
    x=0.72+i*2.45; icon_circle(s, str(i+1), x, 3.0, c); text(s,a,x-0.08,3.9,1.7,0.32,14,WHITE,True); text(s,b,x-0.08,4.3,1.9,0.3,11,MUTED)
    if i<4: arrow(s,x+0.82,3.36,x+2.15,3.36,rgb("526485"),1.4)
card(s, 0.72, 5.15, 3.65, 1.08, "效率", "重复劳动多，批量场景难以扩展", ORANGE, 11)
card(s, 4.82, 5.15, 3.65, 1.08, "质量", "复杂背景与竖排排版容易失真", PINK, 11)
card(s, 8.92, 5.15, 3.65, 1.08, "一致性", "人名、术语、风格难以长期统一", CYAN, 11)

# 3 solution
s = prs.slides.add_slide(blank); add_bg(s, "解决方案：一条面向漫画场景的流水线", "02 / solution", 3)
text(s, "输入一张漫画图，输出可直接阅读、可下载、可继续调整的译图。", 0.66, 1.95, 11, 0.35, 16, WHITE)
flow=[("INPUT","漫画图片",CYAN),("DETECT","文本 / 气泡",PURPLE),("READ","OCR + 语言",PINK),("REBUILD","擦除 + 修复",ORANGE),("WRITE","翻译 + 排版",GREEN),("OUTPUT","预览 / 下载",CYAN)]
for i,(tag,label,c) in enumerate(flow):
    x=0.65+i*2.12; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,3.0,1.72,1.15,fill=PANEL,line=c); text(s,tag,x,3.18,1.72,0.25,10,c,True,MONO,PP_ALIGN.CENTER); text(s,label,x+0.08,3.62,1.56,0.28,13,WHITE,True,PP_ALIGN.CENTER)
    if i<5: arrow(s,x+1.75,3.58,x+2.05,3.58,c,1.6)
text(s, "工程重点：每个阶段都有可替换引擎、状态反馈与失败回退。", 0.7, 5.05, 10.8, 0.4, 15, CYAN, True)
pill(s, "单图任务", 0.72, 5.85, 1.25, CYAN); pill(s, "批量编排", 2.15, 5.85, 1.3, PURPLE); pill(s, "术语一致", 3.65, 5.85, 1.3, PINK); pill(s, "结果可追踪", 5.15, 5.85, 1.55, GREEN)

# 4 architecture
s = prs.slides.add_slide(blank); add_bg(s, "系统总体架构", "03 / architecture", 4)
layers=[("前端交互层", "React 18 · Vite · Tailwind v4 · Framer Motion · GSAP", CYAN), ("业务服务层", "FastAPI · API Router · TranslationTaskManager · SQLite / 文件存储", PURPLE), ("算法引擎层", "CTD / DBNet · MIT48 · manga-ocr · PaddleOCR · Translator · Renderer", PINK), ("图像与模型层", "OpenCV / Pillow · TELEA · LaMa · PyTorch · GPU / CPU 回退", ORANGE)]
for i,(a,b,c) in enumerate(layers):
    y=1.95+i*1.15; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,1.0,y,11.25,0.82,fill=PANEL,line=c); text(s,a,1.3,y+0.18,2.15,0.3,15,c,True); text(s,b,3.7,y+0.18,7.95,0.3,13,WHITE); text(s, f"L{i+1}",11.55,y+0.2,0.45,0.25,9,MUTED,False,MONO,PP_ALIGN.RIGHT)
    if i<3: arrow(s,6.62,y+0.83,6.62,y+1.05,rgb("546481"),1.2)
note(s, "真实代码依据：backend/app/api、services/pipeline.py、services/engines、frontend/src")

# 5 pipeline
s = prs.slides.add_slide(blank); add_bg(s, "核心 Pipeline：先恢复画面，再组织翻译", "04 / pipeline", 5)
stages=[("detect","文本检测"),("ocr","OCR 识别"),("inpaint","原文擦除 / 背景修复"),("group","气泡分组"),("translate","上下文翻译"),("quality","质量检查"),("render","横竖排渲染")]
for i,(tag,label) in enumerate(stages):
    x=0.5+i*1.82; c=[CYAN,PURPLE,ORANGE,PINK,CYAN,GREEN,PURPLE][i]
    icon_circle(s, str(i+1), x+0.48, 2.35, c, 0.58); text(s,tag,x,3.1,1.55,0.2,9,c,True,MONO,PP_ALIGN.CENTER); text(s,label,x,3.45,1.55,0.58,12,WHITE,True,PP_ALIGN.CENTER)
    if i<6: arrow(s,x+1.1,2.64,x+1.68,2.64,rgb("526485"),1.4)
card(s, 0.72, 4.65, 3.75, 1.18, "为什么先修复？", "气泡分组在清理图上执行，原图轮廓作为边界证据，避免相邻气泡被误合并。", ORANGE, 11)
card(s, 4.8, 4.65, 3.75, 1.18, "为什么按气泡翻译？", "整块上下文一次翻译，保留语序与语气，再回填到区域。", PINK, 11)
card(s, 8.88, 4.65, 3.75, 1.18, "为什么可解释？", "每阶段记录进度、耗时、后端与回退原因，便于调试与演示。", CYAN, 11)

# 6 OCR
s = prs.slides.add_slide(blank); add_bg(s, "漫画 OCR：不仅是“读出文字”", "05 / ocr", 6)
image_frame(s, JP, 0.68, 1.9, 4.35, 3.95, "真实日文样例")
text(s, "漫画场景难点", 5.55, 1.98, 2.5, 0.3, 15, WHITE, True)
items=[("01", "竖排与横排并存", CYAN),("02", "艺术字、网点与复杂背景", PURPLE),("03", "检测框、掩膜与置信度要协同", PINK),("04", "拟声词 / 刊头等非气泡文字需保护", ORANGE)]
for i,(n,v,c) in enumerate(items):
    y=2.55+i*0.72; text(s,n,5.58,y,0.45,0.25,10,c,True,MONO); text(s,v,6.2,y,5.6,0.28,14,WHITE)
text(s, "MIT48 + manga-ocr 混合路线：manga-ocr 只补救 MIT 完全未识别的空行。", 5.58, 5.55, 6.5, 0.45, 12, CYAN, True)
note(s, "样例来自 test_image_japanese；图片仅用于展示真实输入，不虚构识别准确率")

# 7 translation
s = prs.slides.add_slide(blank); add_bg(s, "AI 翻译与专有名词：让上下文保持一致", "06 / translation", 7)
for i,(tag,label,c) in enumerate([("OCR","原文区域",CYAN),("CTX","气泡上下文",PURPLE),("TERM","术语匹配",PINK),("LLM","翻译 / 润色",ORANGE),("CHECK","质量告警",GREEN)]):
    x=0.7+i*2.43; icon_circle(s,tag,x,2.35,c,0.82); text(s,label,x-0.2,3.35,1.25,0.3,13,WHITE,True,PP_ALIGN.CENTER)
    if i<4: arrow(s,x+0.9,2.76,x+2.02,2.76,rgb("526485"),1.5)
card(s, 0.72, 4.55, 3.65, 1.35, "六种语言方向", "中文 / 日语 / 英语之间支持全方向配置；默认自动识别 → 中文。", CYAN, 11)
card(s, 4.82, 4.55, 3.65, 1.35, "翻译后端回退", "DeepSeek、OpenAI、DeepL、Google、MyMemory；记录实际使用后端与失败原因。", PINK, 11)
card(s, 8.92, 4.55, 3.65, 1.35, "术语库", "SQLite 存储，支持新增、编辑、删除、搜索与 JSON 导入；整词替换保持人名统一。", PURPLE, 11)

# 8 inpaint
s = prs.slides.add_slide(blank); add_bg(s, "背景修复：OpenCV 与 LaMa 的双路径", "07 / image restoration", 8)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.8,2.0,5.35,3.75,fill=PANEL,line=CYAN); text(s,"CV 路线",1.15,2.35,1.4,0.35,18,CYAN,True); text(s,"适合：平坦、简单背景",1.15,2.95,3.8,0.3,14,WHITE); text(s,"局部背景重建 + TELEA\n速度快 · 资源消耗低\n无额外神经网络权重",1.15,3.55,3.8,1.05,15,MUTED,valign=MSO_ANCHOR.TOP)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,7.15,2.0,5.35,3.75,fill=PANEL,line=PINK); text(s,"LaMa 路线",7.5,2.35,1.7,0.35,18,PINK,True); text(s,"适合：纹理、人物、复杂背景",7.5,2.95,4.2,0.3,14,WHITE); text(s,"局部原尺寸推理\n修复质量更高 · 需 PyTorch / 权重\n可通过配置切换",7.5,3.55,4.1,1.05,15,MUTED,valign=MSO_ANCHOR.TOP)
arrow(s,6.18,3.9,7.08,3.9,GREEN,2.3); text(s,"按配置选择",5.85,3.42,1.45,0.25,10,GREEN,True,MONO,PP_ALIGN.CENTER)
note(s, "当前代码提供 cv / lama 两种 inpainter backend；“自动按场景选择”仍属于后续增强，不在已实现清单中")

# 9 GPU
s = prs.slides.add_slide(blank); add_bg(s, "GPU 加速：优先使用，失败可回退", "08 / runtime", 9)
text(s, "GPU 是可配置的运行时能力，不把未经实测的倍数写成性能承诺。",0.68,1.92,11,0.35,16,WHITE)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.85,2.65,5.2,2.3,fill=PANEL,line=CYAN); icon_circle(s,"GPU",1.25,3.25,CYAN,0.85); text(s,"请求设备",2.35,3.0,2.5,0.3,15,WHITE,True); text(s,"PaddleOCR: gpu:0\nMIT / LaMa: auto / cuda / cpu",2.35,3.55,2.9,0.65,13,MUTED,valign=MSO_ANCHOR.TOP)
arrow(s,6.1,3.8,7.12,3.8,GREEN,2.2)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,7.2,2.65,5.2,2.3,fill=PANEL,line=GREEN); icon_circle(s,"CPU",7.62,3.25,GREEN,0.85); text(s,"安全回退",8.72,3.0,2.5,0.3,15,WHITE,True); text(s,"GPU 初始化失败 / 显存不足\n记录 fallback reason，继续完成任务",8.72,3.55,3.0,0.65,13,MUTED,valign=MSO_ANCHOR.TOP)
pill(s,"串行推理：保护共享模型状态",4.3,5.75,4.7,PURPLE)

# 10 layout
s = prs.slides.add_slide(blank); add_bg(s, "智能排版：翻译完成只是中间结果", "09 / layout", 10)
image_frame(s, EN, 0.72, 1.95, 3.45, 4.15, "真实英文样例")
text(s,"排版约束",4.7,2.02,2.0,0.3,16,WHITE,True)
constraints=[("方向", "中文 / 日语支持横竖排自适应"),("空间", "以气泡几何为边界，避免越界覆盖人物"),("换行", "英文优先按词间换行，中文做均衡分行"),("回退", "气泡掩膜 → 安全扩展框 → 紧致文本框")]
for i,(a,b) in enumerate(constraints):
    y=2.65+i*0.7; pill(s,a,4.72,y,0.82,[CYAN,PINK,ORANGE,GREEN][i]); text(s,b,5.8,y+0.02,6.1,0.3,14,WHITE)
text(s,"排版目标：译文看起来像原生漫画文字，而不是贴在图片上的字幕。",4.72,5.72,7.0,0.4,15,CYAN,True)

# 11 UI
s = prs.slides.add_slide(blank); add_bg(s, "产品闭环：从上传到结果预览", "10 / product loop", 11)
text(s,"已落地的第一阶段体验：任务流式界面 + 底部输入区 + 单张结果卡片。",0.68,1.9,11.5,0.35,16,WHITE)
ui=[("选择 / 粘贴图片",CYAN),("提交任务",PURPLE),("实时进度",PINK),("立即预览",GREEN),("下载 / 重试",ORANGE)]
for i,(lab,c) in enumerate(ui):
    x=0.8+i*2.45; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.85,1.85,1.6,fill=PANEL,line=c); icon_circle(s,str(i+1),x+0.64,3.08,c,0.55); text(s,lab,x+0.12,3.82,1.6,0.3,13,WHITE,True,PP_ALIGN.CENTER)
    if i<4: arrow(s,x+1.9,3.65,x+2.3,3.65,rgb("526485"),1.5)
card(s,0.82,5.15,3.6,0.95,"任务状态","每张图片独立状态、进度、预览与下载",CYAN,11)
card(s,4.86,5.15,3.6,0.95,"批次管理","批量任务汇总进度，可导出 ZIP",PURPLE,11)
card(s,8.9,5.15,3.6,0.95,"术语管理","翻译工具与专有名词页面同级切换",PINK,11)

# 12 real assets / demos
s = prs.slides.add_slide(blank); add_bg(s, "真实样例：输入素材与系统输出", "11 / evidence", 12)
image_frame(s, JP, 0.65, 1.85, 3.7, 4.5, "日语输入")
image_frame(s, EN, 4.82, 1.85, 3.7, 4.5, "英语输入")
image_frame(s, ZH, 8.99, 1.85, 3.7, 4.5, "中文样例")
note(s,"素材来自仓库 test_image_japanese / test_image_english / test_image_chinese；此页展示真实项目素材，不虚构指标")

# 13 engineering challenges
s = prs.slides.add_slide(blank); add_bg(s, "关键工程难点与对应策略", "12 / engineering", 13)
ch=[("气泡误合并", "清理图寻找内部区域 + 原图长轮廓否决跨气泡合并", CYAN), ("擦除残影", "多边形、Otsu 笔画、检测掩膜、注音扩展取并集", ORANGE), ("竖排译文拥挤", "按阅读顺序分列，中文列宽与字距均衡", PINK), ("模型共享状态", "TranslationTaskManager 单 worker + pipeline lock 串行推理", PURPLE), ("翻译不稳定", "多后端回退 + 质量告警 + 实际后端记录", GREEN), ("非气泡文字", "刊头 / 拟声词 / 拉丁标签保护，原文保留", CYAN)]
for i,(a,b,c) in enumerate(ch):
    col=i%2; row=i//2; x=0.72+col*6.1; y=1.95+row*1.47
    card(s,x,y,5.55,1.12,a,b,c,11)

# 14 achieved vs future
s = prs.slides.add_slide(blank); add_bg(s, "实现边界：已完成什么，下一步做什么", "13 / scope", 14)
card(s,0.72,1.92,5.75,4.25,"已实现", "• 三语六方向翻译配置\n• 单图与批量任务编排、进度、ZIP 导出\n• MIT / Paddle OCR 路线与语言路由\n• CV / LaMa 修复引擎切换\n• 气泡分组、质量告警、横竖排渲染\n• 专有名词库 CRUD + JSON 导入\n• 结果预览、下载、失败重试", GREEN, 14)
card(s,6.85,1.92,5.75,4.25,"后续规划", "• 人工文本框拖拽 / 缩放 / 新增编辑器\n• 历史任务与章节级可恢复队列\n• EPUB / CBZ / PDF 输入\n• 更大规模术语知识库与团队协作\n• 自动按场景选择 CV / LaMa\n• 云端 GPU、企业级任务调度\n• 更系统的质量评估与模型融合", PURPLE, 14)
note(s,"已实现项依据 README、project_status、frontend/src、backend/app 代码交叉核对；未实现项明确归入规划")

# 15 summary
s = prs.slides.add_slide(blank); add_bg(s, "项目价值：三层能力形成完整闭环", "14 / conclusion", 15)
vals=[("自动识别","把漫画中的区域、语言与方向变成结构化数据",CYAN),("智能翻译","结合上下文、术语库与多后端回退",PINK),("漫画级重建","擦除、修复、排版后输出可读译图",PURPLE)]
for i,(a,b,c) in enumerate(vals):
    x=0.85+i*4.15; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.2,3.45,2.2,fill=PANEL,line=c); icon_circle(s,str(i+1),x+1.35,2.55,c,0.72); text(s,a,x+0.35,3.45,2.75,0.32,18,WHITE,True,PP_ALIGN.CENTER); text(s,b,x+0.35,3.95,2.75,0.48,12,MUTED,False,PP_ALIGN.CENTER)
arrow(s,4.38,3.3,4.9,3.3,rgb("526485"),1.4); arrow(s,8.52,3.3,9.05,3.3,rgb("526485"),1.4)
text(s,"不是调用一个翻译 API，而是面向漫画场景的可解释图像处理系统。",1.0,5.45,11.4,0.45,19,CYAN,True,PP_ALIGN.CENTER)

# 16 Q&A
s = prs.slides.add_slide(blank); add_bg(s, "THANK YOU", "15 / q&a", 16)
fit_picture(s, LOGO, 5.7, 1.55, 1.95, 1.95)
text(s,"Q&A",3.1,4.0,7.1,0.8,40,WHITE,True,MONO,PP_ALIGN.CENTER)
text(s,"AI Manga Translator",3.1,4.95,7.1,0.35,16,CYAN,True,MONO,PP_ALIGN.CENTER)
text(s,"自动识别 · 智能翻译 · 漫画级图像重建",2.7,5.52,7.9,0.3,12,MUTED,False,PP_ALIGN.CENTER)

prs.core_properties.title = "AI Manga Translator 软件工程项目答辩"
prs.core_properties.subject = "基于 OCR、AI 翻译与图像修复的智能漫画翻译系统"
prs.core_properties.author = ""
prs.save(OUT)
print(OUT)
