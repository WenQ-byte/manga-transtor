from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "AI_Manga_Translator_软件工程项目答辩.pptx"
W,H=13.333,7.5
BG=RGBColor(9,15,27); PANEL=RGBColor(20,29,48); WHITE=RGBColor(241,246,252)
MUTED=RGBColor(158,174,201); CYAN=RGBColor(40,220,240); PURPLE=RGBColor(173,111,255)
PINK=RGBColor(244,104,198); ORANGE=RGBColor(255,178,89); GREEN=RGBColor(80,220,170)
FONT="Microsoft YaHei"; MONO="Consolas"
prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H); blank=prs.slide_layouts[6]

def rgb(h): return RGBColor.from_string(h.replace('#','').upper())
def sh(slide,t,x,y,w,h,fill=None,line=None):
    s=slide.shapes.add_shape(t,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(1)
    return s
def tx(slide,v,x,y,w,h,size=16,color=WHITE,bold=False,font=FONT,align=PP_ALIGN.LEFT,valign=MSO_ANCHOR.MIDDLE):
    if not isinstance(font,str): align,font=font,FONT
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); f=b.text_frame; f.clear(); f.word_wrap=True
    f.margin_left=f.margin_right=Inches(.05); f.margin_top=f.margin_bottom=Inches(.03); f.vertical_anchor=valign
    p=f.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=v; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return b
def bg(slide,title,kicker,n):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=BG
    for x in [.6,3.6,6.6,9.6,12.6]: sh(slide,MSO_SHAPE.RECTANGLE,x,.6,.012,6.3,rgb('18233A'))
    sh(slide,MSO_SHAPE.RECTANGLE,0,0,W,.06,CYAN); tx(slide,"AI MANGA TRANSLATOR  /  DEFENSE",.6,.25,6,.2,8,MUTED,False,MONO); tx(slide,f"{n:02d} / 10",11.9,.25,.8,.2,9,MUTED,False,MONO,PP_ALIGN.RIGHT)
    tx(slide,kicker.upper(),.65,.95,5,.25,10,CYAN,True,MONO); tx(slide,title,.62,1.22,12,.5,27,WHITE,True)
def pill(slide,v,x,y,w,c=CYAN): sh(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,.34,rgb('16283C'),c); tx(slide,v,x,y+.01,w,.28,9,c,True,MONO,PP_ALIGN.CENTER)
def card(slide,x,y,w,h,title,body,c=CYAN,bs=12):
    sh(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,PANEL,rgb('2A3A59')); sh(slide,MSO_SHAPE.RECTANGLE,x,y,.06,h,c)
    tx(slide,title,x+.22,y+.15,w-.4,.3,14,WHITE,True); tx(slide,body,x+.22,y+.55,w-.4,h-.67,bs,MUTED,valign=MSO_ANCHOR.TOP)
def pic(slide,path,x,y,w,h,label):
    sh(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,PANEL,rgb('304262'))
    with Image.open(path) as im: iw,ih=im.size
    r=w/h; ir=iw/ih
    if ir>r: cw=int(ih*r); left=(iw-cw)//2; top=0; ch=ih
    else: ch=int(iw/r); top=(ih-ch)//2; left=0; cw=iw
    p=slide.shapes.add_picture(str(path),Inches(x+.08),Inches(y+.08),Inches(w-.16),Inches(h-.16)); p.crop_left=left/iw; p.crop_right=(iw-left-cw)/iw; p.crop_top=top/ih; p.crop_bottom=(ih-top-ch)/ih
    pill(slide,label,x+.18,y+.16,max(1.1,.12*len(label)+.35),PINK)
def arrow(slide,x1,y1,x2,y2,c=CYAN):
    l=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2)); l.line.color.rgb=c; l.line.width=Pt(2); l.line.end_arrowhead=True

JP=ROOT/'test_image_japanese'/'e19eca8c467765ba30042abd639840e6.jpg'; EN=ROOT/'test_image_english'/'屏幕截图 2026-08-31 114913.png'; ZH=ROOT/'test_image_chinese'/'屏幕截图 2026-08-31 145655.png'; LOGO=ROOT/'frontend'/'public'/'comicx-app-icon.png'

# 1
s=prs.slides.add_slide(blank); bg(s,'AI Manga Translator','software engineering project defense',1)
tx(s,'基于 OCR、AI 翻译与图像修复的智能漫画翻译系统',.7,2.05,8,.55,22,WHITE,True); tx(s,'把识别、翻译、擦除、修复、排版串成一条可运行的工程化流水线',.72,2.78,7.2,.35,13,MUTED); pill(s,'React / FastAPI / MIT OCR / OpenCV / LaMa',.72,3.45,4.15,CYAN)
pic(s,LOGO,9.5,1.55,2.55,2.55,'PROJECT'); tx(s,'姓名 / 小组：____________\n日期：2026',8.95,5.05,3.5,.65,13,MUTED)

# 2
s=prs.slides.add_slide(blank); bg(s,'我们解决什么问题？','01 / problem → value',2)
tx(s,'传统漫画翻译需要在 OCR、翻译、修图、排版之间反复切换；复杂背景和竖排文字尤其耗时。',.7,1.9,11.8,.4,16,WHITE)
for i,(a,b,c) in enumerate([('识别','人工抄录',ORANGE),('翻译','逐句处理',PINK),('修复','擦除 + 补背景',PURPLE),('排版','逐框调整',CYAN)]):
    x=.9+i*3.0; sh(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,3.0,2.05,1.2,PANEL,c); tx(s,a,x,3.2,2.05,.28,16,c,True,PP_ALIGN.CENTER); tx(s,b,x,3.65,2.05,.25,12,MUTED,False,PP_ALIGN.CENTER)
    if i<3: arrow(s,x+2.15,3.6,x+2.75,3.6,rgb('526485'))
tx(s,'目标：输入漫画图片，输出可以直接阅读、预览、下载，并允许继续修正的中文译图。',.9,5.25,11.3,.45,18,CYAN,True,PP_ALIGN.CENTER)

# 3 stack
s=prs.slides.add_slide(blank); bg(s,'项目技术栈：前后端分离 + 可插拔引擎','02 / stack',3)
rows=[('语言','Python  /  JavaScript',CYAN),('前端','React 18 · Vite 6 · Tailwind CSS v4 · Framer Motion · GSAP · Lucide',PURPLE),('后端','FastAPI · Pydantic Settings · SQLite · 文件存储',PINK),('AI / 图像','CTD / DBNet · MIT48 · manga-ocr · PaddleOCR · OpenCV · LaMa · PyTorch',ORANGE)]
for i,(a,b,c) in enumerate(rows):
    y=2.0+i*.95; sh(s,MSO_SHAPE.ROUNDED_RECTANGLE,.85,y,11.65,.65,PANEL,rgb('2A3A59')); tx(s,a,1.15,y+.17,1.55,.25,14,c,True); tx(s,b,2.85,y+.17,9.1,.25,13,WHITE)
tx(s,'工程选择：模型推理保持串行；引擎通过工厂与缓存复用，便于切换与回退。',.95,6.15,11.3,.3,13,MUTED,False,MONO,PP_ALIGN.CENTER)

# 4 flow
s=prs.slides.add_slide(blank); bg(s,'核心流程：翻译不是唯一难点','03 / workflow',4)
st=[('1','检测'),('2','OCR'),('3','擦除 / 修复'),('4','气泡分组'),('5','上下文翻译'),('6','质量检查'),('7','排版输出')]
for i,(n,lbl) in enumerate(st):
    x=.45+i*1.84; c=[CYAN,PURPLE,ORANGE,PINK,CYAN,GREEN,PURPLE][i]; sh(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.75,1.55,1.05,PANEL,c); tx(s,n,x,2.9,1.55,.25,10,c,True,MONO,PP_ALIGN.CENTER); tx(s,lbl,x,3.25,1.55,.25,12,WHITE,True,PP_ALIGN.CENTER)
    if i<6: arrow(s,x+1.58,3.27,x+1.8,3.27,rgb('526485'))
card(s,.85,4.75,3.65,1.05,'亮点 01','在清理图上分组，同时用原图轮廓防止气泡误合并。',CYAN,11); card(s,4.85,4.75,3.65,1.05,'亮点 02','按气泡翻译，保留上下文、语气与阅读顺序。',PINK,11); card(s,8.85,4.75,3.65,1.05,'亮点 03','每张图片独立任务，支持进度、预览、下载、重试。',GREEN,11)

# 5 highlights
s=prs.slides.add_slide(blank); bg(s,'项目亮点与创新点','04 / highlights',5)
for i,(a,b,c) in enumerate([('漫画级 OCR','竖排、艺术字、非气泡文字保护',CYAN),('双路径修复','CV / LaMa 按配置切换，复杂背景有更高质量选项',ORANGE),('智能排版','中文横竖排自适应，气泡几何约束 + 多级回退',PINK),('可解释回退','翻译后端、GPU/CPU、模型失败原因均可追踪',PURPLE),('术语一致','角色名、作品名等术语支持 CRUD 与 JSON 导入',GREEN),('批量闭环','最多 100 张图片编排，结果汇总并导出 ZIP',CYAN)]):
    x=.75+(i%3)*4.15; y=2.05+(i//3)*1.65; card(s,x,y,3.65,1.25,a,b,c,11)

# 6 ai collaboration
s=prs.slides.add_slide(blank); bg(s,'AI 协作：我负责判断，AI 负责执行','05 / ai collaboration',6)
tx(s,'评分关注的不只是“用了 AI”，而是能否提出约束、验证结果、纠正问题并持续迭代。',.72,1.85,11.6,.35,15,WHITE)
turns=[('我提出约束','“先完整阅读代码，所有功能以真实实现为准。”',CYAN),('AI 分析方案','梳理架构、功能边界、素材与可展示证据。',PURPLE),('我验证并追问','发现“规划能力”不能写成“已实现”，要求拆分清单。',PINK),('AI 迭代交付','生成 PPT、核对 PPTX、补充大纲与真实性说明。',GREEN)]
for i,(a,b,c) in enumerate(turns):
    y=2.45+i*.72; sh(s,MSO_SHAPE.ROUNDED_RECTANGLE,.95,y,11.25,.52,PANEL,rgb('2A3A59')); pill(s,a,1.15,y+.09,1.35,c); tx(s,b,2.8,y+.11,8.9,.28,12,WHITE)
tx(s,'对话页说明：以上为真实协作过程的结构化摘录；正式展示时可切换到 OpenCode 的原始聊天截图。',.9,5.75,11.5,.35,11,MUTED,False,MONO,PP_ALIGN.CENTER)

# 7 demo
s=prs.slides.add_slide(blank); bg(s,'项目展示：3–4 分钟完成一次闭环','06 / live demo',7)
pic(s,JP,.7,1.9,3.45,3.9,'输入：日文漫画'); pic(s,EN,4.95,1.9,3.45,3.9,'输入：英文漫画'); pic(s,ZH,9.2,1.9,3.45,3.9,'输出：中文样例')
tx(s,'建议现场演示顺序',.85,6.1,2.2,.25,13,CYAN,True); tx(s,'上传 / 粘贴 → 选择语言 → 提交 → 看任务进度 → 预览结果 → 下载或重试',3.0,6.1,9.4,.25,13,WHITE)

# 8 product
s=prs.slides.add_slide(blank); bg(s,'产品闭环与批量能力','07 / product',8)
for i,(a,b,c) in enumerate([('输入','选择、拖入或粘贴图片',CYAN),('处理','任务状态与阶段进度',PURPLE),('结果','图片预览、下载、重试',GREEN),('批量','多图任务汇总 + ZIP',ORANGE)]):
    x=.8+i*3.05; sh(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.4,2.45,1.55,PANEL,c); tx(s,a,x,2.72,2.45,.3,18,c,True,PP_ALIGN.CENTER); tx(s,b,x+.15,3.3,2.15,.3,12,WHITE,False,PP_ALIGN.CENTER)
    if i<3: arrow(s,x+2.5,3.18,x+2.93,3.18,rgb('526485'))
card(s,1.0,4.85,5.25,1.0,'专有名词管理','独立页面；支持搜索、新增、编辑、删除、JSON 导入。',PINK,12); card(s,7.05,4.85,5.25,1.0,'录屏建议','耗时步骤播放 frontend/public/屏幕录制 2026-09-01 143634.mp4。',CYAN,12)

# 9 boundary
s=prs.slides.add_slide(blank); bg(s,'实现边界：用事实回答老师追问','08 / scope',9)
card(s,.8,1.95,5.7,4.15,'已实现','三语六方向 · 单图/批量任务 · OCR 与语言路由\nCV / LaMa 修复 · 气泡分组 · 质量告警\n横竖排渲染 · 术语库 · 预览下载 · 失败重试',GREEN,15)
card(s,6.85,1.95,5.7,4.15,'后续规划','人工文本框拖拽/缩放/新增\n历史任务与章节级可恢复队列\nEPUB / CBZ / PDF · 团队协作\n云端 GPU · 自动按场景选择修复引擎',PURPLE,15)
tx(s,'回答原则：已经跑通的展示效果讲清楚；尚未完成的功能主动说明边界。',1.05,6.35,11.2,.3,14,CYAN,True,PP_ALIGN.CENTER)

# 10
s=prs.slides.add_slide(blank); bg(s,'总结 / Q&A','09 / close',10)
tx(s,'自动识别',1.0,2.1,3.2,.45,24,CYAN,True,PP_ALIGN.CENTER); tx(s,'结构化文本区域、语言与方向',1.0,2.75,3.2,.35,12,MUTED,False,PP_ALIGN.CENTER)
tx(s,'智能翻译',5.05,2.1,3.2,.45,24,PINK,True,PP_ALIGN.CENTER); tx(s,'上下文 + 术语 + 多后端回退',5.05,2.75,3.2,.35,12,MUTED,False,PP_ALIGN.CENTER)
tx(s,'漫画级重建',9.1,2.1,3.2,.45,24,PURPLE,True,PP_ALIGN.CENTER); tx(s,'擦除、修复、排版后输出译图',9.1,2.75,3.2,.35,12,MUTED,False,PP_ALIGN.CENTER)
tx(s,'不是调用一个翻译 API，而是做了一条面向漫画场景的可解释图像处理流水线。',1.05,4.35,11.2,.5,19,WHITE,True,PP_ALIGN.CENTER); tx(s,'谢谢 / Q&A',4.0,5.55,5.3,.6,30,CYAN,True,MONO,PP_ALIGN.CENTER)

prs.core_properties.title='AI Manga Translator 软件工程项目答辩'; prs.core_properties.subject='7-8分钟项目展示版'; prs.save(OUT); print(OUT)
